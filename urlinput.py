#スクレイピング＋統計＋重回帰
import os
import requests
from bs4 import BeautifulSoup
import time
from datetime import datetime
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt
import glob
import shutil
import scipy.stats as stats
from scipy.stats import linregress
import statsmodels.api as sm
from sklearn.model_selection import train_test_split
from sklearn.linear_model import LinearRegression
from sklearn.metrics import mean_squared_error
from statsmodels.stats.outliers_influence import variance_inflation_factor

#---書き換えあり＜調査駅（基本ローマ字）、URL,page数₍1ページ;100件）>---
#サイトの都合で50以下必須！(多すぎなく3～20程度推奨;case by case)
num_pages = 3

# フォルダーのパス
folder_path = r"C:\\1111accommodation"

# フォルダーを作成（すでに存在する場合はスルー）
os.makedirs(folder_path, exist_ok=True)
# ここでスクリプトを実行

# 調査駅たち
locations = ["品川", "立会川", "梅屋敷", "京急鶴見", "生麦"]
urls = [
        "https://suumo.jp/jj/chintai/ichiran/FR301FC005/?ar=030&bs=040&ra=013&rn=0095&ek=009517460&cb=0.0&ct=9999999&mb=0&mt=9999999&et=9999999&cn=9999999&shkr1=03&shkr2=03&shkr3=03&shkr4=03&sngz=&po1=25&po2=99&pc=100&page=",
        "https://suumo.jp/jj/chintai/ichiran/FR301FC005/?ar=030&bs=040&ra=013&rn=0095&ek=009523090&cb=0.0&ct=9999999&mb=0&mt=9999999&et=9999999&cn=9999999&shkr1=03&shkr2=03&shkr3=03&shkr4=03&sngz=&po1=25&po2=99&pc=100&page=",
        "https://suumo.jp/jj/chintai/ichiran/FR301FC005/?ar=030&bs=040&ra=013&rn=0095&ek=009504660&cb=0.0&ct=9999999&mb=0&mt=9999999&et=9999999&cn=9999999&shkr1=03&shkr2=03&shkr3=03&shkr4=03&sngz=&po1=25&po2=99&pc=100&page=",
        "https://suumo.jp/jj/chintai/ichiran/FR301FC005/?ar=030&bs=040&ra=014&rn=0095&ek=009513460&cb=0.0&ct=9999999&mb=0&mt=9999999&et=9999999&cn=9999999&shkr1=03&shkr2=03&shkr3=03&shkr4=03&sngz=&po1=25&po2=99&pc=100&page=",
        "https://suumo.jp/jj/chintai/ichiran/FR301FC005/?ar=030&bs=040&ra=014&rn=0095&ek=009527900&cb=0.0&ct=9999999&mb=0&mt=9999999&et=9999999&cn=9999999&shkr1=03&shkr2=03&shkr3=03&shkr4=03&sngz=&po1=25&po2=99&pc=100&page="
]

# `zip()` を使って location と url をペアにしてループ
for location, base_url in zip(locations, urls):

    # 空のリストを用意
    all_dataframes = []
    
    # 各ページをスクレイピング
    for i in range(1, num_pages + 1):
        url = base_url + str(i)
        print(f"📡 取得中: {url}")
        
        time.sleep(0.5)  # 0.5秒待機 サイト負荷軽減
        
        response = requests.get(url)
        if response.status_code != 200:
            print(f"⚠️ ページ {i} の取得失敗")
            continue
        
        soup = BeautifulSoup(response.text, "html.parser")
    
        # 物件名・URLの取得
        titles = [title.text.strip() for title in soup.find_all("h2", class_="property_inner-title")]
        links = [a["href"] for a in soup.find_all("a", href=True) if "/chintai/bc" in a["href"]]
        full_links = ["https://suumo.jp" + link for link in links][:100]
    
        # 賃料の取得（数値変換）
        prices = [title.text.strip() for title in soup.find_all("div", class_="detailbox-property-point")]
        
        # 価格を変換する関数
        def convert_price(price):
            return int(float(price.replace('万円', '')) * 10000)
        
        # 価格を変換してリストに格納
        rents = [convert_price(price) for price in prices]
    
        #徒歩時間をだす
        walk_times = []
        # 徒歩時間を含む要素を特定（class="font-weight:bold" または style="font-weight:bold"）
        detail_notes = soup.find_all("div", class_="font-weight:bold") + soup.find_all("div", style="font-weight:bold")
        for note in detail_notes:
            text = note.text.strip()
            try:
                if "歩" in text and "分" in text and "バス" not in text and "車" not in text:
                    # "歩" と "分" の間の数字を抽出
                    walk_time_str = text.split("歩")[1].split("分")[0].strip()
                    walk_time = int(walk_time_str)
                    walk_times.append(walk_time)
                else:
                    walk_times.append(None)
            except (ValueError, IndexError) as e:
                print(f"⚠️ 変換できないデータ: {text}, エラー: {e}")
                walk_times.append(None)
        
        # **物件情報 DataFrame**
        df1 = pd.DataFrame({
            "物件名": titles,
            "URL": full_links,
            "賃料(円)": rents,
            "徒歩時間(分)": walk_times
        })
    
        # **詳細情報 DataFrame**
        properties = []
        for row in soup.find_all("tr")[:100]:  # 100件まで取得
            property_data = {
                "管理費": row.find("td", class_="detailbox-property-col detailbox-property--col1").find_all("div")[1].text.strip(),
                "敷金": row.find("td", class_="detailbox-property-col detailbox-property--col2").find_all("div")[0].text.strip(),
                "礼金": row.find("td", class_="detailbox-property-col detailbox-property--col2").find_all("div")[1].text.strip(),
                "間取り": row.find("td", class_="detailbox-property-col detailbox-property--col3").find_all("div")[0].text.strip(),
                "専有面積(㎡)": row.find("td", class_="detailbox-property-col detailbox-property--col3").find_all("div")[1].text.strip(),
                "向き": row.find("td", class_="detailbox-property-col detailbox-property--col3").find_all("div")[2].text.strip(),
                "物件種別": row.find_all("td", class_="detailbox-property-col detailbox-property--col3")[1].find_all("div")[0].text.strip(),
                "築年数(年)": row.find_all("td", class_="detailbox-property-col detailbox-property--col3")[1].find_all("div")[1].text.strip(),
                "住所": row.find_all("td", class_="detailbox-property-col")[-1].text.strip()
            }
            properties.append(property_data)
    
        df2 = pd.DataFrame(properties)
    
        # **専有面積をfloat型に変換**
        df2["専有面積(㎡)"] = df2["専有面積(㎡)"].str.replace("m2", "").astype(float)
    
        # **築年数をint型に変換**
        df2["築年数(年)"] = pd.to_numeric(df2["築年数(年)"].str.replace("築", "").str.replace("年", "").str.replace("新築", "0"), errors="coerce").astype("Int64")
        df2["築年数(年)"] = df2["築年数(年)"].fillna(0).astype(int)
    
        # **データフレームの結合（横方向）**
        df_combined = pd.concat([df1, df2], axis=1)
    
        # **結合したDataFrameをリストに追加**
        all_dataframes.append(df_combined)
    
    # **すべてのページを行方向（縦方向）に結合**
    final_df = pd.concat(all_dataframes, ignore_index=True)
    df_sorted = final_df.sort_values(by="物件名", ascending=True)
    # 「物件名」列に「築」が含まれる行を削除
    df_sorted = df_sorted[~df_sorted["物件名"].str.contains("築", na=False)]
    # 「物件名」列に「号室」が含まれる行を削除
    df_sorted = df_sorted[~df_sorted["物件名"].str.contains("号室", na=False)]
    #向きが書いていない物件は同一物件率が高い
    df_sorted = df_sorted[~df_sorted["向き"].str.contains("-", na=False)]
    # 対象列を指定
    columns = ['賃料(円)', '管理費', '間取り', '専有面積(㎡)', '向き']
    # 最終行から逆順に、一つ上の行と値が同じ行を特定して削除
    df_sorted = df_sorted.loc[~df_sorted[columns].eq(df_sorted[columns].shift(-1)).all(axis=1)]
    # インデックスをリセットして繰り上げる
    df_sorted = df_sorted.reset_index(drop=True)
    #もう一回間取りでソートしたのち物件名でソートしなおし、対象列が全て同じなら削除する(同一物件が多すぎるため)
    df_sorted = df_sorted.sort_values(by="間取り", ascending=True)
    df_sorted = df_sorted.sort_values(by="物件名", ascending=True)
    # 対象列を指定
    columns = ['賃料(円)', '管理費', '間取り', '専有面積(㎡)', '向き']
    # 最終行から逆順に、一つ上の行と値が同じ行を特定して削除
    df_sorted = df_sorted.loc[~df_sorted[columns].eq(df_sorted[columns].shift(-1)).all(axis=1)]
    #更に一回向きでソートしたのち物件名でソートしなおし、対象列が全て同じなら削除する(同一物件が多すぎるため)
    df_sorted = df_sorted.sort_values(by="向き", ascending=True)
    df_sorted = df_sorted.sort_values(by="物件名", ascending=True)
    # 対象列を指定
    columns = ['賃料(円)', '管理費', '間取り', '専有面積(㎡)', '向き']
    # 最終行から逆順に、一つ上の行と値が同じ行を特定して削除
    df_sorted = df_sorted.loc[~df_sorted[columns].eq(df_sorted[columns].shift(-1)).all(axis=1)]
    #更に一回賃料でソートしたのち物件名でソートしなおし、対象列が全て同じなら削除する(同一物件が多すぎるため)
    df_sorted = df_sorted.sort_values(by="賃料(円)", ascending=True)
    df_sorted = df_sorted.sort_values(by="物件名", ascending=True)
    # 対象列を指定
    columns = ['賃料(円)', '管理費', '間取り', '専有面積(㎡)', '向き']
    # 最終行から逆順に、一つ上の行と値が同じ行を特定して削除
    df_sorted = df_sorted.loc[~df_sorted[columns].eq(df_sorted[columns].shift(-1)).all(axis=1)]
    # インデックスをリセットして繰り上げる
    df_sorted = df_sorted.reset_index(drop=True)
    
    #行数をだす
    n=len(df_sorted)
    
    # 現在時刻を取得し、"yyMMddhhmm" 形式に変換
    timestamp = datetime.now().strftime("%y%m%d%H%M")
    datestamp = datetime.now().strftime("%y%m%d")
    
    # ファイル名を動的に生成
    file_name = "1fData_"+f"{location}_{datestamp}.csv"
    
    # フォルダのパス（例: Cドライブ accommodation フォルダ）
    full_path = os.path.join(folder_path, file_name)
    
    #NAを削除する
    df_sorted = df_sorted.replace('', pd.NA).dropna()
    
    # CSVファイルを保存
    df_sorted.to_csv(full_path, index=False, encoding="utf-8-sig")
    
    print(f"n=",n)
    
    #----基礎統計----
    
    # 現在時刻を取得し、変数に保存（固定）
    current_time = datetime.now()
    #対象の駅を表示
    print(location)
    
    df_base1 = np.array([
        ["全データ数", "取得した現在時刻", "調査駅", "出典"],
        [n, "day"+timestamp, location, "https://suumo.jp/jj/chintai"]
    ])
    
    
    print(df_base1)
    
    #df_base1:基礎情報の整備
    #賃料(円)の統計データ
    #df_base1:基礎情報の整備
    #賃料(円)の統計データ
    avg_total_rents = round(df_sorted["賃料(円)"].mean(),2) #平均
    medi_total_rents = round(df_sorted["賃料(円)"].median(),2) #中央値
    stdevs_total_rents = round(df_sorted["賃料(円)"].std(ddof=1),4) #不変標準偏差
    std_error_total_rents =  round(df_sorted["賃料(円)"].std(ddof=1)/np.sqrt(len(df_sorted)),4) #標準誤差
    min_total_rents = df_sorted["賃料(円)"].min() #最小値
    max_total_rents = df_sorted["賃料(円)"].max() #最大値
    firstQ_total_rents = round(df_sorted["賃料(円)"].quantile(0.25),1) #第一四分位   
    thirdQ_total_rents = round(df_sorted["賃料(円)"].quantile(0.75),1) #第三四分位
    kurt_total_rents = round(df_sorted["賃料(円)"].kurt(),2) #尖度  
    skew_total_rents = round(df_sorted["賃料(円)"].skew(),2) #歪度 
    
    #徒歩時間(分)の統計データ
    # 数値型をすべて float に変換
    df_sorted["徒歩時間(分)"] = df_sorted["徒歩時間(分)"].astype(float)
    
    # 徒歩時間(分)の統計データ
    avg_total_walk_times = round(df_sorted["徒歩時間(分)"].mean(), 2)  # 平均
    medi_total_walk_times = round(df_sorted["徒歩時間(分)"].median(), 2)  # 中央値
    stdevs_total_walk_times = round(df_sorted["徒歩時間(分)"].std(ddof=1), 4)  # 不偏標準偏差
    std_error_total_walk_times = round(df_sorted["徒歩時間(分)"].std(ddof=1) / np.sqrt(len(df_sorted)), 4)  # 標準誤差
    
    # **整数型の影響を受けないよう変換**
    min_total_walk_times = df_sorted["徒歩時間(分)"].min()  # 最小値
    max_total_walk_times = df_sorted["徒歩時間(分)"].max()  # 最大値
    firstQ_total_walk_times = round(df_sorted["徒歩時間(分)"].quantile(0.25), 2)  # 第一四分位
    thirdQ_total_walk_times = round(df_sorted["徒歩時間(分)"].quantile(0.75), 2)  # 第三四分位
    kurt_total_walk_times = round(df_sorted["徒歩時間(分)"].kurt(), 2)  # 尖度
    skew_total_walk_times = round(df_sorted["徒歩時間(分)"].skew(), 2)  # 歪度
    
    #専有面積(㎡)の統計データ
    avg_total_space = round(df_sorted["専有面積(㎡)"].mean(),2) #平均
    medi_total_space = round(df_sorted["専有面積(㎡)"].median(),2) #中央値
    stdevs_total_space = round(df_sorted["専有面積(㎡)"].std(ddof=1),4) #不変標準偏差
    std_error_total_space =  round(df_sorted["専有面積(㎡)"].std(ddof=1)/np.sqrt(len(df_sorted)),4) #標準誤差
    min_total_space = df_sorted["専有面積(㎡)"].min() #最小値
    max_total_space = df_sorted["専有面積(㎡)"].max() #最大値
    firstQ_total_space = round(df_sorted["専有面積(㎡)"].quantile(0.25),2) #第一四分位   
    thirdQ_total_space = round(df_sorted["専有面積(㎡)"].quantile(0.75),2) #第三四分位
    kurt_total_space = round(df_sorted["専有面積(㎡)"].kurt(),2) #尖度  
    skew_total_space = round(df_sorted["専有面積(㎡)"].skew(),2) #歪度 
    
    #築年数(年)の統計データ
    avg_total_ages = round(df_sorted["築年数(年)"].mean(),2) #平均
    medi_total_ages = round(df_sorted["築年数(年)"].median(),2) #中央値
    stdevs_total_ages = round(df_sorted["築年数(年)"].std(ddof=1),4) #不変標準偏差
    std_error_total_ages =  round(df_sorted["築年数(年)"].std(ddof=1)/np.sqrt(len(df_sorted)),4) #標準誤差
    min_total_ages = df_sorted["築年数(年)"].min() #最小値
    max_total_ages = df_sorted["築年数(年)"].max() #最大値
    firstQ_total_ages = round(df_sorted["専有面積(㎡)"].quantile(0.25),1) #第一四分位   
    thirdQ_total_ages = round(df_sorted["築年数(年)"].quantile(0.75),1) #第三四分位
    kurt_total_ages = round(df_sorted["築年数(年)"].kurt(),2) #尖度  
    skew_total_ages = round(df_sorted["築年数(年)"].skew(),2) #歪度  
    
    df_stats1 = np.array([["項目","平均","中央値","不変標準偏差","標準誤差","最小値","最大値","第一四分位","第三四分位","尖度","歪度"],
         ["賃料(円)",avg_total_rents,medi_total_rents,stdevs_total_rents,std_error_total_rents,min_total_rents,max_total_rents,firstQ_total_rents,thirdQ_total_rents,kurt_total_rents,skew_total_rents],
        ["徒歩時間(分)",avg_total_walk_times,medi_total_walk_times,stdevs_total_walk_times,std_error_total_walk_times,min_total_walk_times,max_total_walk_times,firstQ_total_walk_times,thirdQ_total_walk_times,kurt_total_walk_times,skew_total_walk_times],
        ["専有面積(㎡)",avg_total_space,medi_total_space,stdevs_total_space,std_error_total_space,max_total_space,min_total_space,firstQ_total_space,thirdQ_total_space,kurt_total_space,skew_total_space],
        ["築年数(年)",avg_total_ages,medi_total_ages,stdevs_total_ages,std_error_total_ages,min_total_ages,max_total_ages,firstQ_total_ages,thirdQ_total_ages,kurt_total_ages,skew_total_ages]])
    
    #DataFrameを転置する（書きにくかったので）
    df_stats1 = df_stats1.T
    print(df_stats1)
    
    df_stats11 = np.array([["項目","平均","中央値","不変標準偏差","標準誤差"],
         ["賃料(円)",avg_total_rents,medi_total_rents,stdevs_total_rents,std_error_total_rents],
        ["徒歩時間(分)",avg_total_walk_times,medi_total_walk_times,stdevs_total_walk_times,std_error_total_walk_times],
        ["専有面積(㎡)",avg_total_space,medi_total_space,stdevs_total_space,std_error_total_space],
        ["築年数(年)",avg_total_ages,medi_total_ages,stdevs_total_ages,std_error_total_ages]])
    
    #DataFrameを転置する（書きにくかったので）
    df_stats11 = df_stats11.T
    
    df_stats12 = np.array([["項目","最小値","最大値","第一四分位","第三四分位","尖度","歪度"],
         ["賃料(円)",min_total_rents,max_total_rents,firstQ_total_rents,thirdQ_total_rents,kurt_total_rents,skew_total_rents],
        ["徒歩時間(分)",min_total_walk_times,max_total_walk_times,firstQ_total_walk_times,thirdQ_total_walk_times,kurt_total_walk_times,skew_total_walk_times],
        ["専有面積(㎡)",min_total_space,max_total_space,firstQ_total_space,thirdQ_total_space,kurt_total_space,skew_total_space],
        ["築年数(年)",min_total_ages,max_total_ages,firstQ_total_ages,thirdQ_total_ages,kurt_total_ages,skew_total_ages]])
    
    #DataFrameを転置する（書きにくかったので）
    df_stats12 = df_stats12.T
    
    # ファイル名を動的に生成
    # 基本情報と基礎統計量を`DataFrame` に変換してcsv化
    df_base1 = pd.DataFrame(df_base1)
    df_stats1 = pd.DataFrame(df_stats1)
    
    file_name_base1 = f"{location}_{datestamp}_base1.csv"
    file_name_stats1 = f"{location}_{datestamp}_stats1.csv"
    
    full_path_base1 = os.path.join(folder_path, file_name_base1)
    full_path_stats1 = os.path.join(folder_path, file_name_stats1)
    
    df_base1.to_csv(full_path_base1, index=False, encoding="utf-8-sig")
    df_stats1.to_csv(full_path_stats1, index=False, encoding="utf-8-sig")
    
    #DataFrame分割化（大きすぎる）
    df_stats11 = pd.DataFrame(df_stats11)
    df_stats12 = pd.DataFrame(df_stats12)
    
    file_name_stats11 = f"{location}_{datestamp}_stats11.csv"
    file_name_stats12 = f"{location}_{datestamp}_stats12.csv"
    
    full_path_stats11 = os.path.join(folder_path, file_name_stats11)
    full_path_stats12 = os.path.join(folder_path, file_name_stats12)
    
    df_stats11.to_csv(full_path_stats11, index=False, encoding="utf-8-sig")
    df_stats12.to_csv(full_path_stats12, index=False, encoding="utf-8-sig")
    
    #matplotlibを日本語化
    # フォントを変更（WindowsならMS Gothic、MacならHiragino）
    plt.rcParams['font.family'] = 'MS Gothic'  # Windowsの場合
    
    # 図の準備（2行2列のレイアウト）
    fig, axes = plt.subplots(4, 2, figsize=(12, 12))  
    
    # ① 賃料（円）のヒストグラム
    axes[0, 0].hist(df_sorted["賃料(円)"], bins=30, edgecolor='black')
    axes[0, 0].set_title("賃料（円）のヒストグラム")
    axes[0, 0].set_xlabel("賃料（円")
    axes[0, 0].set_ylabel("度数")
    
    # ② 賃料（円）の箱ひげ図
    axes[0, 1].boxplot(df_sorted["賃料(円)"], patch_artist=True, boxprops=dict(facecolor="skyblue"))
    axes[0, 1].set_title("賃料（円）の箱ひげ図")
    axes[0, 1].set_xlabel(location)
    axes[0, 1].set_ylabel("賃料（円")
    
    # ③ 徒歩時間（分）のヒストグラム
    axes[1, 0].hist(df_sorted["徒歩時間(分)"], bins=30, edgecolor='black')
    axes[1, 0].set_title("徒歩時間（分）のヒストグラム")
    axes[1, 0].set_xlabel("徒歩時間（分）")
    axes[1, 0].set_ylabel("度数")
    
    # ④ 専有面積(㎡)の箱ひげ図
    axes[1, 1].boxplot(df_sorted["徒歩時間(分)"], patch_artist=True, boxprops=dict(facecolor="lightgreen"))
    axes[1, 1].set_title("徒歩時間（分）の箱ひげ図")
    axes[1, 1].set_xlabel(location)
    axes[1, 1].set_ylabel("徒歩時間（分")
    
    # ③ 専有面積(㎡)のヒストグラム
    axes[2, 0].hist(df_sorted["専有面積(㎡)"], bins=30, edgecolor='black')
    axes[2, 0].set_title("専有面積(㎡)のヒストグラム")
    axes[2, 0].set_xlabel("専有面積(㎡)")
    axes[2, 0].set_ylabel("度数")
    
    # ④ 専有面積(㎡)の箱ひげ図
    axes[2, 1].boxplot(df_sorted["専有面積(㎡)"], patch_artist=True, boxprops=dict(facecolor="lightgreen"))
    axes[2, 1].set_title("専有面積(㎡)の箱ひげ図")
    axes[2, 1].set_xlabel(location)
    axes[2, 1].set_ylabel("専有面積(㎡)")
    
    # ③ 築年数(年のヒストグラム
    axes[3, 0].hist(df_sorted["築年数(年)"], bins=30, edgecolor='black')
    axes[3, 0].set_title("築年数(年)のヒストグラム")
    axes[3, 0].set_xlabel("築年数(年)")
    axes[3, 0].set_ylabel("度数")
    
    # ④ 築年数(年の箱ひげ図
    axes[3 ,1].boxplot(df_sorted["築年数(年)"], patch_artist=True, boxprops=dict(facecolor="lightgreen"))
    axes[3 ,1].set_title("築年数(年)の箱ひげ図")
    axes[3, 1].set_xlabel(location)
    axes[3 ,1].set_ylabel("築年数(年)")
    
    # 画像として保存
    plt.tight_layout()
    file_name_g1 = f"{location}_{datestamp}_tg1.png"
    full_path_g1 = os.path.join(folder_path, file_name_g1)
    plt.savefig(full_path_g1)  # PNGファイルとして保存
    plt.close(fig)  # 画像を閉じる（メモリ解放のため）
    
    #分布グラフ拡大化
    #賃料
    # 図の準備（2行2列のレイアウト）
    fig, axes = plt.subplots(1, 2, figsize=(12, 6))  # 横並びの画像を作成
    # ① 賃料（円）のヒストグラム
    axes[0].hist(df_sorted["賃料(円)"], bins=30, edgecolor='black')
    axes[0].set_title("賃料(円)のヒストグラム")
    axes[0].set_xlabel("賃料(円)")  # 括弧を修正
    axes[0].set_ylabel("度数")
    # ② 賃料（円）の箱ひげ図
    axes[1].boxplot(df_sorted["賃料(円)"], patch_artist=True, boxprops=dict(facecolor="skyblue"))
    axes[1].set_title("賃料(円)の箱ひげ図")
    axes[1].set_xlabel(location)
    axes[1].set_ylabel("賃料(円)")  # 括弧を修正
    # 画像として保存
    plt.tight_layout()
    file_name_gr1 = f"{location}_{datestamp}_gr1.png"
    full_path_gr1 = os.path.join(folder_path, file_name_gr1)
    plt.savefig(full_path_gr1)  # PNGファイルとして保存
    plt.close(fig)  # 画像を閉じる（メモリ解放のため）

    #徒歩時間
    fig, axes = plt.subplots(1, 2, figsize=(12, 6))  # 横並びの画像を作成
    # 徒歩時間のヒストグラム
    axes[0].hist(df_sorted["徒歩時間(分)"], bins=30, edgecolor='black')
    axes[0].set_title("徒歩時間(分)のヒストグラム")
    axes[0].set_xlabel("徒歩時間(分)")  # 括弧を修正
    axes[0].set_ylabel("度数")
    # 徒歩時間の箱ひげ図
    axes[1].boxplot(df_sorted["徒歩時間(分)"], patch_artist=True, boxprops=dict(facecolor="skyblue"))
    axes[1].set_title("徒歩時間(分)の箱ひげ図")
    axes[1].set_xlabel(location)
    axes[1].set_ylabel("徒歩時間(分)")  # 括弧を修正
    # 画像として保存
    plt.tight_layout()
    file_name_gw1 = f"{location}_{datestamp}_gw1.png"
    full_path_gw1 = os.path.join(folder_path, file_name_gw1)
    plt.savefig(full_path_gw1)  # PNGファイルとして保存
    plt.close(fig)  # 画像を閉じる（メモリ解放のため）

    #専有面積
    fig, axes = plt.subplots(1, 2, figsize=(12, 6))  # 横並びの画像を作成
    # 徒歩時間のヒストグラム
    axes[0].hist(df_sorted["専有面積(㎡)"], bins=30, edgecolor='black')
    axes[0].set_title("専有面積のヒストグラム")
    axes[0].set_xlabel("専有面積(㎡)")  # 括弧を修正
    axes[0].set_ylabel("度数")
    # 専有面積の箱ひげ図
    axes[1].boxplot(df_sorted["専有面積(㎡)"], patch_artist=True, boxprops=dict(facecolor="skyblue"))
    axes[1].set_title("専有面積の箱ひげ図")
    axes[1].set_xlabel(location)
    axes[1].set_ylabel("専有面積(㎡)")  # 括弧を修正
    # 画像として保存
    plt.tight_layout()
    file_name_gs1 = f"{location}_{datestamp}_gs1.png"
    full_path_gs1 = os.path.join(folder_path, file_name_gs1)
    plt.savefig(full_path_gs1)  # PNGファイルとして保存
    plt.close(fig)  # 画像を閉じる（メモリ解放のため）
    
    #築年数
    fig, axes = plt.subplots(1, 2, figsize=(12, 6))  # 横並びの画像を作成
    # 築年数のヒストグラム
    axes[0].hist(df_sorted["築年数(年)"], bins=30, edgecolor='black')
    axes[0].set_title("築年数のヒストグラム")
    axes[0].set_xlabel("築年数(年)")  # 括弧を修正
    axes[0].set_ylabel("度数")
    # 築年数の箱ひげ図
    axes[1].boxplot(df_sorted["築年数(年)"], patch_artist=True, boxprops=dict(facecolor="skyblue"))
    axes[1].set_title("築年数の箱ひげ図")
    axes[1].set_xlabel(location)
    axes[1].set_ylabel("築年数(年)")  # 括弧を修正
    # 画像として保存
    plt.tight_layout()
    file_name_ga1 = f"{location}_{datestamp}_ga1.png"
    full_path_ga1 = os.path.join(folder_path, file_name_ga1)
    plt.savefig(full_path_ga1)  # PNGファイルとして保存
    plt.close(fig)  # 画像を閉じる（メモリ解放のため）

    # X軸の共通データ
    x1 = df_sorted["徒歩時間(分)"]
    x2 = df_sorted["専有面積(㎡)"]
    x3 = df_sorted["築年数(年)"]
    
    # Y軸のデータ
    y1 = df_sorted["賃料(円)"]
    y2 = df_sorted["賃料(円)"]
    y3 = df_sorted["賃料(円)"]
    
    # 図の準備（縦に3つ並べる）
    fig, axes = plt.subplots(3, 1, figsize=(6, 18))
    
    # 散布図のタイトルと軸ラベル
    titles = ["賃料(円) vs 徒歩時間(分)", "専有面積(㎡) vs 賃料(円)", "賃料(円) vs 築年数(年)"]
    x_labels = ["徒歩時間(分)", "専有面積(㎡)", "築年数(年)"]
    y_labels = ["賃料(円)",  "賃料(円)", "賃料(円)"]
    x_values = [x1, x2, x3]
    y_values = [y1, y2, y3]
    
    from scipy.stats import linregress
    # グラフ描画
    for i in range(3):
        # 線形回帰を計算
        slope, intercept, r_value, p_value, std_err = linregress(x_values[i], y_values[i])
        line_eq = f"y = {slope:.2f}x + {intercept:.2f}"
    
        # 散布図を描画
        axes[i].scatter(x_values[i], y_values[i], alpha=0.5, color="blue", label="データ")
    
        # 近似直線を描画
        axes[i].plot(x_values[i], slope*x_values[i] + intercept, color="red", label=f"近似直線: {line_eq}")
    
        # 決定係数とp値を表示
        axes[i].text(min(x_values[i]), max(y_values[i]), f"R² = {r_value**2:.2f}\n p値 = {p_value:.4f}", fontsize=10, color="black")
    
        # タイトル・ラベル設定
        axes[i].set_title(titles[i])
        axes[i].set_xlabel(x_labels[i])
        axes[i].set_ylabel(y_labels[i])
        axes[i].legend(loc="lower right")
    
    # レイアウト調整
    plt.tight_layout()
    
    # 画像保存
    file_name_g2 = f"{location}_{datestamp}_tg2.png"
    full_path_g2 = os.path.join(folder_path, file_name_g2)
    plt.savefig(full_path_g2)  # PNGファイルとして保存
    plt.close(fig)  # 画像を閉じる（メモリ解放のため）

    #散布図を3つに分割して一つづつ出す
    for i in range(3):
        # 新しい図を作成
        plt.figure(figsize=(6, 4))
    
        # 線形回帰を計算
        slope, intercept, r_value, p_value, std_err = linregress(x_values[i], y_values[i])
        line_eq = f"y = {slope:.2f}x + {intercept:.2f}"
    
        # 散布図を描画
        plt.scatter(x_values[i], y_values[i], alpha=0.5, color="blue", label="データ")
    
        # 近似直線を描画
        plt.plot(x_values[i], slope*x_values[i] + intercept, color="red", label=f"近似直線: {line_eq}")
    
        # 決定係数とp値を表示
        plt.text(min(x_values[i]), max(y_values[i]), f"R² = {r_value**2:.2f}\n p値 = {p_value:.4f}", fontsize=10, color="black")
    
        # タイトル・ラベル設定
        plt.title(titles[i])
        plt.xlabel(x_labels[i])
        plt.ylabel(y_labels[i])
        plt.legend(loc="lower right")
    
        # 画像の保存
        file_name_tgscat = f"{location}_{datestamp}_tgscat{i+1}.png"  # 各画像に異なる名前を付ける
        full_path_tgscat = os.path.join(folder_path, file_name_tgscat)
        plt.savefig(full_path_tgscat)  # PNGファイルとして保存
        plt.close(fig)  # 画像を閉じる（メモリ解放のため）
    
    # "間取り" のカテゴリー定義
    categories = ["ワンルーム", "1K", "1DK", "1LDK", "2K", "2DK", "2LDK", "3K", "3DK", "3LDK"]
    df_sorted["間取り分類"] = df_sorted["間取り"].apply(lambda x: x if x in categories else "その他")
    
    # グループ化して件数・平均賃料・平均専有面積を集計
    cat1 = df_sorted.groupby("間取り分類").agg(
        件数=("間取り分類", "count"),
        平均賃料=("賃料(円)", "mean"),
        平均専有面積=("専有面積(㎡)", "mean")
    ).reset_index()
    
    # 小数1桁に丸める
    cat1[["平均賃料", "平均専有面積"]] = cat1[["平均賃料", "平均専有面積"]].round(1)
    
    file_name_cat1 = f"{location}_{datestamp}_ct1.csv"
    full_path_cat1 = os.path.join(folder_path, file_name_cat1)
    cat1.to_csv(full_path_cat1, index=False, encoding="utf-8-sig")
    
    # 結果を表示
    print(cat1)
    
#----重回帰-----
 
    # 説明変数と目的変数を定義
    X = df_sorted[['徒歩時間(分)', '築年数(年)', '専有面積(㎡)']]
    y = df_sorted['賃料(円)']
    
    # 定数項を追加
    X = sm.add_constant(X)
    
    # 線形回帰モデルの作成
    model = sm.OLS(y, X).fit()
    
    # モデルの要約を表示
    print("   ")
    print(model.summary())
    print("-----切片を除いてP>|t|が0.05以下だと有意と考える　それ以上なら本来はその係数を除いて重回帰やり直し　関係性があるとは言い切れない-----")
    print("   ")
    adj_r_squared = model.rsquared_adj  # 補正決定係数
    f_stat = model.fvalue  # F 値
    f_p_value = model.f_pvalue  # F の p 値
    intercept_coef = model.params["const"]  # 切片の係数
    coefficients = model.params.drop("const")  # 説明変数の傾き
    p_values = model.pvalues.drop("const")  # 説明変数の p 値
    
    # **日本語形式で表示**
    print("=== 線形回帰モデルの結果 ===")
    print(f"補正決定係数: {adj_r_squared:.4f}")
    print(f"F値: {f_stat:.4f}")
    print(f"Fのp値: {f_p_value:.4f}")
    print(f"切片の係数: {intercept_coef:.4f}\n")
    
    print("各説明変数の傾きと p 値:")
    for var in coefficients.index:
        print(f" - {var}: 傾き = {coefficients[var]:.4f}, p 値 = {p_values[var]:.4f}")
    
    # ファイル名を動的に生成
    # 重回帰基本情報とを`DataFrame` に変換してcsv化
    
    df_mrl1 = np.array([
        ["指標", "値"],  # ヘッダーを追加
        ["補正決定係数", adj_r_squared],
        ["F値", f_stat],
        ["Fのp値", f_p_value]
    ])
    
    df_mrl1 = pd.DataFrame(df_mrl1).T
    
    file_name_mrl1 = f"{location}_{datestamp}_mrl1.csv"
    full_path_mrl1 = os.path.join(folder_path, file_name_mrl1)
    
    df_mrl1.to_csv(full_path_mrl1, index=False, encoding="utf-8-sig")
    
    print(intercept_coef)
    #print( coefficients[0],coefficients[1],p_values[1],coefficients[2],p_values[2],coefficients[3,],p_values[3])
    
    # 切片傾き情報とを`DataFrame` に変換してcsv化
    df_mrl2 = np.array([
        ["item", "coef(切片、傾き)","p値"],  # ヘッダーを追加
        ["切片", intercept_coef,"-"],  # ヘッダーを追加
        ["徒歩時間(分)", coefficients["徒歩時間(分)"],p_values["徒歩時間(分)"]],
        ["築年数(年)", coefficients["築年数(年)"],p_values["築年数(年)"]],
        ["専有面積(㎡)", coefficients["専有面積(㎡)"],p_values["専有面積(㎡)"]]
    ])
    
    df_mrl2 = pd.DataFrame(df_mrl2)
    
    file_name_mrl2 = f"{location}_{datestamp}_mrl2.csv"
    full_path_mrl2 = os.path.join(folder_path, file_name_mrl2)
    
    df_mrl2.to_csv(full_path_mrl2, index=False, encoding="utf-8-sig")
    
    #予測値と実測値の比較
    # **データの準備**
    df_plot = df_sorted.copy()
    df_plot = df_plot.drop_duplicates()  # 重複削除
    df_plot = df_plot.reset_index(drop=True)  # インデックスをリセット
    
    # **予測値を計算**
    X_pred = sm.add_constant(df_plot[['徒歩時間(分)', '築年数(年)', '専有面積(㎡)']])
    df_plot['predicted_rent'] = model.predict(X_pred)
    
    # **残差の標準誤差を計算**
    residuals = df_plot['賃料(円)'] - df_plot['predicted_rent']
    std_residuals = np.std(residuals)
    
    # **予測区間を残差の標準誤差で近似（Rコードと同様）**
    df_plot['upper_bound'] = df_plot['predicted_rent'] + (std_residuals * 1.96)
    df_plot['lower_bound'] = df_plot['predicted_rent'] - (std_residuals * 1.96)
    
    # **プロット用に予測値でソート（スムーズな線のため）**
    df_plot_sorted = df_plot.sort_values('predicted_rent').reset_index(drop=True)
    
    # **決定係数 (R²) の計算**
    r_squared = model.rsquared
    p_values_model = model.pvalues
    
    # **データ数を取得**
    n_samples = len(df_plot)
    
    # **近似式の作成（実家賃をx、予測家賃をy）**
    slope, intercept = np.polyfit(df_plot['賃料(円)'], df_plot['predicted_rent'], 1)
    line_eq = f"y = {slope:.2f}x + {intercept:.2f}"
    
    # **予測区間の幅を計算**
    gap_pred = std_residuals * 1.96  # gap_predを定義
    
    #予測値幅
    print(f"予測区間の幅（±1.96σ): {gap_pred:.1f}")
    
    # **プロットの作成**
    plt.figure(figsize=(12, 8))
    
    # 散布図（実家賃 vs 予測家賃）
    plt.scatter(df_plot['賃料(円)'], df_plot['predicted_rent'], 
               color="blue", alpha=0.6, label="実測値", s=30)
    
    # **スムーズな線を描画するために十分な点を生成**
    x_smooth = np.linspace(df_plot['賃料(円)'].min(), df_plot['賃料(円)'].max(), 100)
    y_smooth = slope * x_smooth + intercept
    
    # 回帰直線（スムーズ）
    plt.plot(x_smooth, y_smooth, "r-", lw=2, label="回帰直線")
    
    # **予測区間線をスムーズに描画**
    # 実家賃に対応する予測値を計算して区間線を描画
    upper_smooth = y_smooth + gap_pred
    lower_smooth = y_smooth - gap_pred
    
    plt.plot(x_smooth, upper_smooth, "k--", lw=1.5, alpha=0.8, label="予測区間上限")
    plt.plot(x_smooth, lower_smooth, "k--", lw=1.5, alpha=0.8, label="予測区間下限")
    
    # **予測区間の塗りつぶし**
    plt.fill_between(x_smooth, lower_smooth, upper_smooth, 
                     color="orange", alpha=0.2, label="予測区間")
    
    # **95%信頼区間も追加（より狭い区間）**
    confidence_interval = std_residuals * 1.96 / np.sqrt(n_samples)
    upper_conf = y_smooth + confidence_interval
    lower_conf = y_smooth - confidence_interval
    
    plt.fill_between(x_smooth, lower_conf, upper_conf, 
                     color="blue", alpha=0.3, label="95% 信頼区間")
    
    # **グラフの詳細設定**
    plt.xlabel("実際の賃料 (円)", fontsize=12)
    plt.ylabel("予測賃料 (円)", fontsize=12)
    plt.title("実際の賃料 vs 予測賃料（信頼区間・予測区間付き）", fontsize=14)
    plt.legend(loc='upper left')
    plt.grid(True, alpha=0.3)
    
    # **統計情報を右下に表示**
    plt.text(0.98, 0.02, 
             f"近似式: {line_eq}\nR² = {r_squared:.3f}\np値 = {p_values_model[1]:.3f}\nn = {n_samples}",
             fontsize=11, verticalalignment="bottom", horizontalalignment="right",
             transform=plt.gca().transAxes,
             bbox=dict(facecolor="white", alpha=0.8, edgecolor="gray"))
    
    # **画像保存**
    file_name_mlrap1 = f"{location}_{datestamp}_mlrap1.png"
    image_path_mlrap1 = os.path.join(folder_path, file_name_mlrap1)
    plt.savefig(image_path_mlrap1, dpi=300, bbox_inches='tight')  # PNGファイルとして保存
    plt.close(fig)  # 画像を閉じる（メモリ解放のため）
    
    # VIFの計算
    print("-----VIFは多重共線性（マルチコ；リニアリティー）の指数で1に近ければ大体OK-----")
    vif_data = pd.DataFrame()
    vif_data["feature"] = X.columns
    vif_data["VIF"] = [variance_inflation_factor(X.values, i) for i in range(X.shape[1])]
    
    # VIFとを`DataFrame` に変換してcsv化
    df_vif1 = np.array([
        ["item", "VIF"],  # ヘッダーを追加
        ["徒歩時間(分)", variance_inflation_factor(X.values, 1)],
        ["築年数(年)", variance_inflation_factor(X.values, 2)],
        ["専有面積(㎡)", variance_inflation_factor(X.values, 3)]
    ])
    
    df_vif1 = pd.DataFrame(df_vif1)
    
    file_name_vif1 = f"{location}_{datestamp}_vif1.csv"
    full_path_vif1 = os.path.join(folder_path, file_name_vif1)
    
    df_vif1.to_csv(full_path_vif1, index=False, encoding="utf-8-sig")
    
    print(df_vif1)
    
    # **面積別の賃料予測値を計算（切片も含める）**
    pred25 = round(intercept_coef + coefficients["専有面積(㎡)"]*25 + coefficients["徒歩時間(分)"]*10 + coefficients["築年数(年)"]*10, 1)
    pred50 = round(intercept_coef + coefficients["専有面積(㎡)"]*50 + coefficients["徒歩時間(分)"]*10 + coefficients["築年数(年)"]*10, 1)
    pred75 = round(intercept_coef + coefficients["専有面積(㎡)"]*75 + coefficients["徒歩時間(分)"]*10 + coefficients["築年数(年)"]*10, 1)
    pred100 = round(intercept_coef + coefficients["専有面積(㎡)"]*100 + coefficients["徒歩時間(分)"]*10 + coefficients["築年数(年)"]*10, 1)
    
    # **DataFrame を作成**
    df_comp1 = pd.DataFrame([
        ["25m²", pred25, round(pred25 - gap_pred, 1), round(pred25 + gap_pred, 1)],
        ["50m²", pred50, round(pred50 - gap_pred, 1), round(pred50 + gap_pred, 1)],
        ["75m²", pred75, round(pred75 - gap_pred, 1), round(pred75 + gap_pred, 1)],
        ["100m²", pred100, round(pred100 - gap_pred, 1), round(pred100 + gap_pred, 1)]
    ], columns=["専有面積", "予測値", "予測下限", "予測上限"])  # ✅ DataFrame の columns を直接指定
    
    # **結果を表示**
    print(df_comp1)
    
    
    # **CSVファイルとして保存**
    file_name_comp1 = f"{location}_{datestamp}_comp1.csv"
    full_path_comp1 = os.path.join(folder_path, file_name_comp1)
    df_comp1.to_csv(full_path_comp1, index=False, encoding="utf-8-sig")   
    
    # **CSVファイルとして保存**
    file_name_comp1 = f"{location}_{datestamp}_comp1.csv"
    full_path_comp1 = os.path.join(folder_path, file_name_comp1)
    df_comp1.to_csv(full_path_comp1, index=False, encoding="utf-8-sig")
    
    
    # **CSVファイルとして保存**
    file_name_comp1 = f"{location}_{datestamp}_comp1.csv"
    full_path_comp1 = os.path.join(folder_path, file_name_comp1)
    df_comp1.to_csv(full_path_comp1, index=False, encoding="utf-8-sig")
    
    #****ここからPowerPoint****
    
    #ファイル名とフォルダー名を指定
    file_name = f"1e_{location}_{timestamp}_ptt1.pptx"
    file_path = os.path.join(folder_path, file_name)
    
    #プレゼンテーションを作る
    pptt1 = Presentation()
    
    # **スライド枚数を取得**
    num_slides = len(pptt1.slides)
    print(f"初期スライド数: {num_slides}")  # 初期値を確認
    
    # **タイトルスライドを追加（必要なら）**
    if num_slides == 0:
        slide_layout1 = pptt1.slide_layouts[0]  # タイトルスライド
        slide1 = pptt1.slides.add_slide(slide_layout1)
    
        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]
    
        title.text = f"{location}駅\n徒歩圏内の賃貸物件の\n調査結果"
        subtitle.text = f"調査時刻: {timestamp}\nデータ件数は{n}です\n ご注意:重複はなるべく排除していますが排除され切れていません"
    
        # **スライド枚数を更新**
        num_slides = len(pptt1.slides)
        print(f"タイトルスライド追加後のスライド数: {num_slides}")  # ここで確認
    
    # **スライドが1枚なら2枚目を追加**
    if num_slides == 1:
        slide_layout2 = pptt1.slide_layouts[6]  # 空白スライド
        slide2 = pptt1.slides.add_slide(slide_layout2)
        print("✅ 2枚目のスライドを追加しました")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")  # ここで確認
    
    # **「基本情報」のテキストボックスを追加**
    text_box2 = slide2.shapes.add_textbox(Cm(0.4), Cm(0.4), Cm(5), Cm(1))
    text_frame2 = text_box2.text_frame
    p2 = text_frame2.add_paragraph()
    p2.text = "基本情報"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0, 0, 0)  # 黒色
    
    # **データフレーム df_base1 を表として挿入**
    # 表の作成 (1.5cm, 2cm の位置に配置)
    rows, cols = df_base1.shape[0] + 1, df_base1.shape[1]
    table = slide2.shapes.add_table(rows, cols, Cm(1.5), Cm(2), Cm(22), Cm(4)).table
    
    # ヘッダー行の設定
    for col_idx, col_name in enumerate(df_base1.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
    
    # データ行の設定
    for row_idx, row in enumerate(df_base1.itertuples(), start=1):
        for col_idx, value in enumerate(row[1:]):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
    
    # **スライドが2枚なら3枚目を追加**
    if num_slides == 2:
        slide_layout3 = pptt1.slide_layouts[6]  # 空白スライド
        slide3 = pptt1.slides.add_slide(slide_layout3)
        print("✅ 3枚目のスライドを追加しました")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")  # ここで確認
    
    # **「カテゴリー情報」のテキストボックスを追加**
    text_box3 = slide3.shapes.add_textbox(Cm(0.4), Cm(0.4), Cm(5), Cm(1))
    text_frame3 = text_box3.text_frame
    p2 = text_frame3.add_paragraph()
    p2.text = "カテゴリー情報"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0, 0, 0)  # 黒色
    
    # **データフレーム df_cat1 を表として挿入**
    # 表の作成 (1.5cm, 2cm の位置に配置)
    rows, cols = cat1.shape[0] + 1, cat1.shape[1]
    table = slide3.shapes.add_table(rows, cols, Cm(1.5), Cm(2), Cm(22), Cm(15)).table
    
    # ヘッダー行の設定
    for col_idx, col_name in enumerate(cat1.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
    
    # データ行の設定
    for row_idx, row in enumerate(cat1.itertuples(), start=1):
        for col_idx, value in enumerate(row[1:]):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
    
    # **スライドが3枚なら4枚目を追加**
    if num_slides == 3:
        slide_layout4 = pptt1.slide_layouts[6]  # 空白スライド
        slide4 = pptt1.slides.add_slide(slide_layout4)  # スライド追加
        print("✅ 4枚目のスライドを追加しました！")
    
        # **スライド枚数を再取得**
        num_slides = len(pptt1.slides)  # 更新
        print(f"🔄 更新後のスライド数: {num_slides}") 
    
    
    # **「基礎統計情報」のテキストボックスを追加**
    text_box4 = slide4.shapes.add_textbox(Cm(0.4), Cm(0.4), Cm(5), Cm(1))
    text_frame4 = text_box4.text_frame
    p3 = text_frame4.add_paragraph()
    p3.text = "基礎統計量情報；すみません今は大きすぎて表に出来てません　次のスライド以降に表示します"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(0, 0, 0)  # 黒色
    
    # **スライド枚数を再確認**
    print(f"変更後のスライド数: {num_slides}")  # ここで確認
    
    p2.font.color.rgb = RGBColor(0, 0, 0)  # 黒色
    
    # テキストを作成する座標
    left = Cm(1)
    top = Cm(3)
    width = Cm(24)
    height = Cm(16)
    
    # df_stats1のデータを取得（テキスト化）
    text_content = df_stats1.to_string(index=False)
    
    # スライドにテキストボックスを追加
    text_box_stats1 = slide4.shapes.add_textbox(left, top, width, height)
    text_frame = text_box_stats1.text_frame
    text_frame.text = text_content  # データをテキストとして入力
    
    # **スライドが4枚なら5枚目を追加**
    if num_slides == 4:
        slide_layout5 = pptt1.slide_layouts[6]  # 空白スライド
        slide5 = pptt1.slides.add_slide(slide_layout5)
        print("✅ 5枚目のスライドを追加しました")
    
        # **スライド枚数を再取得**
        num_slides = len(pptt1.slides)  # 更新
        print(f"🔄 更新後のスライド数: {num_slides}") 
    
    # **「基礎統計情報A」のテキストボックスを追加**
    text_box5 = slide5.shapes.add_textbox(Cm(0.4), Cm(0.4), Cm(5), Cm(1))
    text_frame5 = text_box5.text_frame
    p5 = text_frame5.add_paragraph()
    p5.text = "基礎統計量情報A"
    p5.font.size = Pt(16)
    p5.font.bold = True
    p5.font.color.rgb = RGBColor(0, 0, 0)  # 黒色
    
    # **データフレーム df_stas11 を表として挿入**
    # 表の作成 (1.5cm, 2cm の位置に配置)
    rows, cols = df_stats11.shape[0] + 1, df_stats11.shape[1]
    table = slide5.shapes.add_table(rows, cols, Cm(1.5), Cm(2), Cm(22), Cm(15)).table
    
    # ヘッダー行の設定
    for col_idx, col_name in enumerate(df_stats11.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
    
    # データ行の設定
    for row_idx, row in enumerate(df_stats11.itertuples(), start=1):
        for col_idx, value in enumerate(row[1:]):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
    
    # **スライドが5枚なら6枚目を追加**
    if num_slides == 5:
        slide_layout6 = pptt1.slide_layouts[6]  # 空白スライド
        slide6 = pptt1.slides.add_slide(slide_layout6)  # スライド追加
        print("✅ 6枚目のスライドを追加しました！")
    
        # **スライド枚数を再取得**
        num_slides = len(pptt1.slides)  # 更新
        print(f"🔄 更新後のスライド数: {num_slides}") 
    
    # **「基礎統計情報B」のテキストボックスを追加**
    text_box6 = slide6.shapes.add_textbox(Cm(0.4), Cm(0.4), Cm(5), Cm(1))
    text_frame6 = text_box6.text_frame
    p6 = text_frame6.add_paragraph()
    p6.text = "基礎統計量情報B"
    p6.font.size = Pt(16)
    p6.font.bold = True
    p6.font.color.rgb = RGBColor(0, 0, 0)  # 黒色
    
    # **データフレーム df_stas12 を表として挿入**
    # 表の作成 (1.5cm, 2cm の位置に配置)
    rows, cols = df_stats12.shape[0] + 1, df_stats12.shape[1]
    table = slide6.shapes.add_table(rows, cols, Cm(1.5), Cm(2), Cm(22), Cm(15)).table
    
    # ヘッダー行の設定
    for col_idx, col_name in enumerate(df_stats12.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
    
    # データ行の設定
    for row_idx, row in enumerate(df_stats12.itertuples(), start=1):
        for col_idx, value in enumerate(row[1:]):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
    
    # **スライドが6枚なら7枚目を追加**
    if len(pptt1.slides) == 6:  # スライド数を直接取得
        slide_layout7 = pptt1.slide_layouts[3]  # 最後のレイアウトを取得
        slide7 = pptt1.slides.add_slide(slide_layout7)  # スライドを追加
        print(f"✅ 7枚目のスライドを追加しました！現在のスライド数: {len(pptt1.slides)}")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")  
    
    # スライドのプレースホルダーを取得
    
    # テキストボックスの処理（通常スライドレイアウトによってプレースホルダーのインデックスが変わる）
    text_boxes = [shape for shape in slide7.shapes if shape.has_text_frame]  # テキストフレームのあるシェイプを取得
    
    # 右側のテキストボックスに"分布"、左側に"一次回帰"を設定（配置が違う場合は調整）
    if len(text_boxes) >= 2:
        text_boxes[0].text = "全体の分布と一次回帰のグラフ"  # タイトル
        text_boxes[1].text = "分布"  # 右側
        text_boxes[2].text = "一次回帰"  # 左側
        # 画像ファイルのパスを作成
    
    tg1 = os.path.normpath(os.path.join(folder_path, f"{location}_{datestamp}_tg1.png"))
    tg2 = os.path.normpath(os.path.join(folder_path, f"{location}_{datestamp}_tg2.png"))
    
    if not os.path.exists(tg1):
        print(f"❌ ファイルが見つかりません: {tg1}")
    else:
        print(f"✅ ファイルが見つかりました: {tg1}")
    
    # スライドのプレースホルダーを取得
    left_placeholder = text_boxes[1]  # 左側のプレイスホルダー
    right_placeholder = text_boxes[2]  # 右側のプレイスホルダー
    
    # 画像を追加（プレイスホルダーの位置に合わせる）
    slide7.shapes.add_picture(tg1, left_placeholder.left, left_placeholder.top, left_placeholder.width, left_placeholder.height)
    slide7.shapes.add_picture(tg2, right_placeholder.left, right_placeholder.top, right_placeholder.width, right_placeholder.height)
    
    print("✅ 左側のプレイスホルダーに tg1.png、右側に tg2.png を追加しました！")
    
    # **スライドが7枚なら8枚目を追加**
    if len(pptt1.slides) == 7:  
        slide_layout8 = pptt1.slide_layouts[5]  # レイアウト取得
        slide8 = pptt1.slides.add_slide(slide_layout8)  # スライド追加
        print(f"✅ 8枚目のスライドを追加しました！ 現在のスライド数: {len(pptt1.slides)}")
    
        file_name_gr1 = f"{location}_{datestamp}_gr1.png"
        image_path_gr1 = os.path.normpath(os.path.join(folder_path, file_name_gr1))  # パスを正規化
    
        # **タイトルを設定**
        title8 = slide8.shapes.title
        if title8:  # タイトルが存在するかチェック
            title8.text = "賃料分布グラフ"
            # **画像の位置を調整**
        left = Inches(0.3)    # 左から 1 インチ
        top = Inches(1.5)     # 上から 2 インチ
        width = Inches(9.5)   # 幅 5 インチ
        height = Inches(4.5)  # 高さ 4 インチ
    
        # **画像ファイルの追加**
        if os.path.exists(image_path_gr1):  # 画像ファイルの存在を確認
            slide8.shapes.add_picture(image_path_gr1, left, top, width, height)  # 正しい変数を使用
            print(f"✅ 画像をスライドに追加しました: {image_path_gr1}")
        else:
            print(f"❌ 画像ファイルが見つかりません: {image_path_gr1}")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")
    
    # **スライドが8枚なら9枚目を追加**
    if len(pptt1.slides) == 8:  
        slide_layout9 = pptt1.slide_layouts[5]  # レイアウト取得
        slide9 = pptt1.slides.add_slide(slide_layout9)  # スライド追加
        print(f"✅ 9枚目のスライドを追加しました！ 現在のスライド数: {len(pptt1.slides)}")
    
        file_name_gw1 = f"{location}_{datestamp}_gw1.png"
        image_path_gw1 = os.path.normpath(os.path.join(folder_path, file_name_gw1))  # パスを正規化
    
        # **タイトルを設定**
        title9 = slide9.shapes.title
        if title9:  # タイトルが存在するかチェック
            title9.text = "徒歩時間グラフ"
            # **画像の位置を調整**
        left = Inches(0.3)    # 左から 1 インチ
        top = Inches(1.5)     # 上から 2 インチ
        width = Inches(9.5)   # 幅 5 インチ
        height = Inches(4.5)  # 高さ 4 インチ
    
        # **画像ファイルの追加**
        if os.path.exists(image_path_gw1):  # 画像ファイルの存在を確認
            slide9.shapes.add_picture(image_path_gw1, left, top, width, height)  # 正しい変数を使用
            print(f"✅ 画像をスライドに追加しました: {image_path_gw1}")
        else:
            print(f"❌ 画像ファイルが見つかりません: {image_path_gw1}")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")
    
    # **スライドが9枚なら10枚目を追加**
    if len(pptt1.slides) == 9:  
        slide_layout10 = pptt1.slide_layouts[5]  # レイアウト取得
        slide10 = pptt1.slides.add_slide(slide_layout10)  # スライド追加
        print(f"✅ 枚目のスライドを追加しました！ 現在のスライド数: {len(pptt1.slides)}")
    
        file_name_gs1 = f"{location}_{datestamp}_gs1.png"
        image_path_gs1 = os.path.normpath(os.path.join(folder_path, file_name_gs1))  # パスを正規化
    
        # **タイトルを設定**
        title10 = slide10.shapes.title
        if title10:  # タイトルが存在するかチェック
            title10.text = "専有面積グラフ"
            # **画像の位置を調整**
        left = Inches(0.3)    # 左から 1 インチ
        top = Inches(1.5)     # 上から 2 インチ
        width = Inches(9.5)   # 幅 5 インチ
        height = Inches(4.5)  # 高さ 4 インチ
    
        # **画像ファイルの追加**
        if os.path.exists(image_path_gs1):  # 画像ファイルの存在を確認
            slide10.shapes.add_picture(image_path_gs1, left, top, width, height)  # 正しい変数を使用
            print(f"✅ 画像をスライドに追加しました: {image_path_gs1}")
        else:
            print(f"❌ 画像ファイルが見つかりません: {image_path_gs1}")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")
        
    # **スライドが10枚なら11枚目を追加**
    if len(pptt1.slides) == 10:  
        slide_layout11 = pptt1.slide_layouts[5]  # レイアウト取得
        slide11 = pptt1.slides.add_slide(slide_layout11)  # スライド追加
        print(f"✅ 11枚目のスライドを追加しました！ 現在のスライド数: {len(pptt1.slides)}")
    
        file_name_ga1 = f"{location}_{datestamp}_ga1.png"
        image_path_ga1 = os.path.normpath(os.path.join(folder_path, file_name_ga1))  # パスを正規化
    
        # **タイトルを設定**
        title11 = slide11.shapes.title
        if title11:  # タイトルが存在するかチェック
            title11.text = "築年数グラフ"
            # **画像の位置を調整**
        left = Inches(0.3)    # 左から 1 インチ
        top = Inches(1.5)     # 上から 2 インチ
        width = Inches(9.5)   # 幅 5 インチ
        height = Inches(4.5)  # 高さ 4 インチ
    
        # **画像ファイルの追加**
        if os.path.exists(image_path_ga1):  # 画像ファイルの存在を確認
            slide11.shapes.add_picture(image_path_ga1, left, top, width, height)  # 正しい変数を使用
            print(f"✅ 画像をスライドに追加しました: {image_path_ga1}")
        else:
            print(f"❌ 画像ファイルが見つかりません: {image_path_ga1}")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")
    
    
    # **スライドが11枚なら12枚目を追加**
    if len(pptt1.slides) == 11:  # スライド数を直接取得
        slide_layout12 = pptt1.slide_layouts[5]  # レイアウト取得
        slide12 = pptt1.slides.add_slide(slide_layout12)  # スライド追加
        print(f"✅ 12枚目のスライドを追加しました！ 現在のスライド数: {len(pptt1.slides)}")
    
        file_name_tgscat1 = f"{location}_{datestamp}_tgscat1.png"
        image_path_tgscat1 = os.path.normpath(os.path.join(folder_path, file_name_tgscat1))  # パスを正規化
    
        # **タイトルを設定**
        title12 = slide12.shapes.title
        if title12:  # タイトルが存在するかチェック
            title12.text = "賃料と徒歩時間の散布図"
            # **画像の位置を調整**
        left = Inches(0.3)    # 左から 1 インチ
        top = Inches(1.5)     # 上から 2 インチ
        width = Inches(9.5)   # 幅 5 インチ
        height = Inches(5.0)  # 高さ 4 インチ
    
        # **画像ファイルの追加**
        if os.path.exists(image_path_tgscat1):  # 画像ファイルの存在を確認
            slide12.shapes.add_picture(image_path_tgscat1, left, top, width, height)  # 正しい変数を使用
            print(f"✅ 画像をスライドに追加しました: {image_path_tgscat1}")
        else:
            print(f"❌ 画像ファイルが見つかりません: {image_path_tgscat1}")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")
    
    # **スライドが12枚なら13枚目を追加**
    if len(pptt1.slides) == 12:  # スライド数を直接取得
        slide_layout13 = pptt1.slide_layouts[5]  # レイアウト取得
        slide13 = pptt1.slides.add_slide(slide_layout13)  # スライド追加
        print(f"✅ 13枚目のスライドを追加しました！ 現在のスライド数: {len(pptt1.slides)}")
    
        file_name_tgscat2 = f"{location}_{datestamp}_tgscat2.png"
        image_path_tgscat2 = os.path.normpath(os.path.join(folder_path, file_name_tgscat2))  # パスを正規化
    
        # **タイトルを設定**
        title13 = slide13.shapes.title
        if title13:  # タイトルが存在するかチェック
            title13.text = "賃料と専有面積の散布図"
            # **画像の位置を調整**
        left = Inches(0.3)    # 左から 1 インチ
        top = Inches(1.5)     # 上から 2 インチ
        width = Inches(9.5)   # 幅 5 インチ
        height = Inches(5.0)  # 高さ 4 インチ
    
        # **画像ファイルの追加**
        if os.path.exists(image_path_tgscat1):  # 画像ファイルの存在を確認
            slide13.shapes.add_picture(image_path_tgscat2, left, top, width, height)  # 正しい変数を使用
            print(f"✅ 画像をスライドに追加しました: {image_path_tgscat2}")
        else:
            print(f"❌ 画像ファイルが見つかりません: {image_path_tgscat2}")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")
    
    # **スライドが13枚なら14枚目を追加**
    if len(pptt1.slides) == 13:  # スライド数を直接取得
        slide_layout14 = pptt1.slide_layouts[5]  # レイアウト取得
        slide14 = pptt1.slides.add_slide(slide_layout14)  # スライド追加
        print(f"✅ 14枚目のスライドを追加しました！ 現在のスライド数: {len(pptt1.slides)}")
    
        file_name_tgscat3 = f"{location}_{datestamp}_tgscat3.png"
        image_path_tgscat3 = os.path.normpath(os.path.join(folder_path, file_name_tgscat3))  # パスを正規化
    
        # **タイトルを設定**
        title14 = slide14.shapes.title
        if title14:  # タイトルが存在するかチェック
            title14.text = "賃料と築年数の散布図"
            # **画像の位置を調整**
        left = Inches(0.3)    # 左から 1 インチ
        top = Inches(1.5)     # 上から 2 インチ
        width = Inches(9.5)   # 幅 5 インチ
        height = Inches(5.0)  # 高さ 4 インチ
    
        # **画像ファイルの追加**
        if os.path.exists(image_path_tgscat3):  # 画像ファイルの存在を確認
            slide14.shapes.add_picture(image_path_tgscat3, left, top, width, height)  # 正しい変数を使用
            print(f"✅ 画像をスライドに追加しました: {image_path_tgscat3}")
        else:
            print(f"❌ 画像ファイルが見つかりません: {image_path_tgscat3}")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")
    
    for index, slide in enumerate(pptt1.slides):
        # スライドの番号
        current_page = index + 1  
        total_pages = len(pptt1.slides)  
    
    # **スライドが14枚なら15枚目を追加**
    if num_slides == 14:
        slide_layout15 = pptt1.slide_layouts[5]  # 空白スライド
        slide15 = pptt1.slides.add_slide(slide_layout15)
        print("✅ 15枚目のスライドを追加しました")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")  # ここで確認
    
    # **タイトルとサブタイトルを設定**
    title = slide.shapes.title
    
    
    if slide15.shapes.title:
        slide15.shapes.title.text = "重回帰分析結果"
    else:
        print("⚠ スライドにタイトルプレースホルダーがありません。")
    
    # **「基礎統計情報A」のテキストボックスを追加**
    text_box15 = slide15.shapes.add_textbox(Cm(0.4), Cm(2), Cm(5), Cm(1))
    text_frame15 = text_box15.text_frame
    p15 = text_frame15.add_paragraph()  # `text_frame5` を `text_frame15` に修正
    p15.text = "重回帰基礎結果とcoefficients"
    p15.font.size = Pt(16)
    p15.font.bold = True
    p15.font.color.rgb = RGBColor(0, 0, 0)  # 黒色
    
    # **table15_1 をスライドの上から4cmの位置に配置**
    table15_1 = slide15.shapes.add_table(df_mrl1.shape[0], df_mrl1.shape[1], Cm(1.5), Cm(4.0), Cm(22), Cm(4)).table
    
    # **table15_2 をスライドの上から10cmの位置に配置（高さ調整）**
    table15_2 = slide15.shapes.add_table(df_mrl2.shape[0], df_mrl2.shape[1], Cm(1.5), Cm(10.0), Cm(22), Cm(4)).table
    
    # **table15_1 に df_mrl1 のデータを追加**
    for row_idx, (index, row) in enumerate(df_mrl1.iterrows()):
        for col_idx, value in enumerate(row):
            cell = table15_1.cell(row_idx, col_idx)  # 1行目はヘッダーなので +1
            cell.text = str(value)
    
    # **table15_2 に df_mrl2 のデータを追加**
    for row_idx, (index, row) in enumerate(df_mrl2.iterrows()):
        for col_idx, value in enumerate(row):
            cell = table15_2.cell(row_idx, col_idx)  # 1行目はヘッダーなので +1
            cell.text = str(value)
    
    # **スライドが15枚なら16枚目を追加**
    if num_slides == 15:
        slide_layout16 = pptt1.slide_layouts[6]  # 空白スライド
        slide16 = pptt1.slides.add_slide(slide_layout16)
        print("✅ 16枚目のスライドを追加しました")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")  # ここで確認
    
    # **テキストボックスを追加**
    text_box16 = slide16.shapes.add_textbox(Cm(0.4), Cm(0.5), Cm(5), Cm(1))
    text_frame16 = text_box16.text_frame
    p16 = text_frame16.add_paragraph() 
    p16.text = "重回帰の多重共線性（VIF)"
    p16.font.size = Pt(16)
    p16.font.bold = True
    p16.font.color.rgb = RGBColor(0, 0, 0)  # 黒色
    
    # **データフレーム df_vif1 を表として挿入**
    # 表の作成 (1.5cm, 2cm の位置に配置)
    rows, cols = df_vif1.shape[0] + 1, df_vif1.shape[1]
    table = slide16.shapes.add_table(rows, cols, Cm(1.5), Cm(2.5), Cm(22), Cm(4)).table
    
    # ヘッダー行の設定
    for col_idx, col_name in enumerate(df_vif1.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
    
    # データ行の設定
    for row_idx, row in enumerate(df_vif1.itertuples(), start=1):
        for col_idx, value in enumerate(row[1:]):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
    
    # **スライドが16枚なら17枚目を追加**
    if num_slides == 16:
        slide_layout17 = pptt1.slide_layouts[5]  # 空白スライド
        slide17 = pptt1.slides.add_slide(slide_layout17)
        print("✅ 17枚目のスライドを追加しました")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")  # ここで確認
    
        # **タイトルを設定**
        title17 = slide17.shapes.title
        if title17:  # タイトルが存在するかチェック
            title17.text = "予測家賃と実家賃の関係"
            # **画像の位置を調整**
        left = Inches(0.3)    # 左から 1 インチ
        top = Inches(1.5)     # 上から 2 インチ
        width = Inches(9.5)   # 幅 5 インチ
        height = Inches(5.0)  # 高さ 4 インチ
    
        # **画像ファイルの追加**
        if os.path.exists(image_path_mlrap1):  # 画像ファイルの存在を確認
            slide17.shapes.add_picture(image_path_mlrap1, left, top, width, height)  # 正しい変数を使用
            print(f"✅ 画像をスライドに追加しました: {image_path_mlrap1}")
        else:
            print(f"❌ 画像ファイルが見つかりません: {image_path_mlrap1}")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")
    
    # **スライドが17枚なら18枚目を追加**
    if num_slides == 17:
        slide_layout18 = pptt1.slide_layouts[6]  # 空白スライド
        slide18 = pptt1.slides.add_slide(slide_layout18)
        print("✅ 18枚目のスライドを追加しました")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")  # ここで確認
    
    # **テキストボックスを追加**
    text_box18 = slide18.shapes.add_textbox(Cm(0.4), Cm(0.5), Cm(5), Cm(1))
    text_frame18 = text_box18.text_frame
    p18 = text_frame18.add_paragraph() 
    p18.text = "面積毎の家賃予測"
    p18.font.size = Pt(16)
    p18.font.bold = True
    p18.font.color.rgb = RGBColor(0, 0, 0)  # 黒色
    
    # **データフレーム df_comp1 を表として挿入**
    # 表の作成 (1.5cm, 2cm の位置に配置)
    rows, cols = df_comp1.shape[0] + 1, df_comp1.shape[1]
    table = slide18.shapes.add_table(rows, cols, Cm(1.5), Cm(2), Cm(22), Cm(4)).table
    
    # ヘッダー行の設定
    for col_idx, col_name in enumerate(df_comp1.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
    
    # データ行の設定
    for row_idx, row in enumerate(df_comp1.itertuples(), start=1):
        for col_idx, value in enumerate(row[1:]):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
    
    # **スライドが18枚なら19枚目を追加**
    if num_slides == 18:
        slide_layout19 = pptt1.slide_layouts[6]  # 空白スライド
        slide19 = pptt1.slides.add_slide(slide_layout19)
        print("✅ 19枚目のスライドを追加しました")
    
        # **スライド枚数を再確認**
        num_slides = len(pptt1.slides)
        print(f"変更後のスライド数: {num_slides}")  # ここで確認
    
    # フッター用のテキスト
    for index, slide in enumerate(pptt1.slides):
        # スライドの番号
        current_page = index + 1  
        total_pages = len(pptt1.slides)  
    
        # フッター用のテキスト
        left_text = f"{location}, n={n}"
        center_text = f"{current_page}/{total_pages}"
        right_text = f"{timestamp}"
    
        # **テキストボックスの追加（スライド下部）**
        left_box = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(2), Inches(0.3))
        left_box.text_frame.text = left_text  
    
        center_box = slide.shapes.add_textbox(Inches(4.2), Inches(7.15), Inches(2), Inches(0.3))
        center_box.text_frame.text = center_text  
    
        right_box = slide.shapes.add_textbox(Inches(8.0), Inches(7.15), Inches(2), Inches(0.3))
        right_box.text_frame.text = right_text  
        left_text = f"{location}, n={n}"
        center_text = f"{current_page}/{total_pages}"
        right_text = f"{timestamp}"
    
        # **テキストボックスの追加（スライド下部）**
        left_box = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(2), Inches(0.3))
        left_box.text_frame.text = left_text  
    
        center_box = slide.shapes.add_textbox(Inches(4.2), Inches(7.15), Inches(2), Inches(0.3))
        center_box.text_frame.text = center_text  
    
        right_box = slide.shapes.add_textbox(Inches(8.0), Inches(7.15), Inches(2), Inches(0.3))
        right_box.text_frame.text = right_text  
    
    
    # **PowerPointを保存**
    pptt1.save(file_path)
    print(f"✅ ファイルを保存しました: {file_path}")

#総合まとめpptxを作成
# フォルダー内のファイルを取得
csv_files = [f for f in os.listdir(folder_path) if f.startswith("1fData") and f.endswith(".csv")]

# データを格納するリスト
data_list = []

# CSVファイルを処理
for file in csv_files:
    file_path = os.path.join(folder_path, file)
    
    # ファイル名から「_ と _ の間の文字列」を抽出
    name_parts = file.split("_")
    if len(name_parts) >= 3:  # 少なくとも2つの `_` があることを確認
        column_name = name_parts[1]  # 2番目の要素を列名として使用

        # CSVを読み込む
        df = pd.read_csv(file_path)

        # 「賃料（円）」の列を抽出（正確な列名に合わせて調整）
        rent_column = [col for col in df.columns if "賃料" in col or "円" in col]  # 「賃料（円）」を探す
        if rent_column:
            df_filtered = df[[rent_column[0]]]  # 最初に見つかった「賃料」列を取得
            df_filtered.columns = [column_name]  # 抽出した列の名前を変更
            data_list.append(df_filtered)

# 複数のデータを結合
if data_list:
    result_df = pd.concat(data_list, axis=1)
    print(result_df)
else:
    print("該当するデータが見つかりませんでした。")

# 基礎統計量を求める
stats_df = result_df.describe()
print("基礎統計量:\n", stats_df)

# 📌 【箱ひげ図の作成・保存】
plt.figure(figsize=(10, 6))
result_df.boxplot()
plt.title(f"箱ひげ図 ({datestamp})")  # タイトルに datestamp を追加
plt.ylabel("賃料（円）")
plt.xticks(rotation=45)
plt.grid(True)

filename_box1 = f"{datestamp}_box1.png"  # ファイル名に datestamp を追加
image_path_box1 = os.path.join(folder_path, filename_box1)
plt.savefig(image_path_box1)  # 正しいパスで保存
plt.close(fig)  # 画像を閉じる（メモリ解放のため）

# 列名を取得
column_names = result_df.columns.tolist()  # `result_df` のすべての列名をリスト化

# ANOVAの実施（動的な列を利用）
# 列名を取得
column_names = result_df.columns.tolist()

# 欠損値を削除してデータを整理
result_df = result_df.dropna()

# ANOVAの実施
groups = [result_df[col] for col in column_names]

# データが空でないか確認
groups = [g for g in groups if len(g) > 0]

if len(groups) > 1:  # 2つ以上のグループがある場合に適用
    F_value, p_value = stats.f_oneway(*groups)

    # 結果をテキストにまとめる
    text_anova1 = f"一元配置分散分析（ANOVA）の結果:\nF値: {F_value:.2f}\np値: {p_value:.3f}"
    
    # 出力
    print(text_anova1)
else:
    print("ANOVAを適用できる十分なデータがありません。")


# 📌 【賃料の累積比率グラフの作成・保存】
plt.figure(figsize=(10, 6))

for column in result_df.columns:
    data = np.sort(result_df[column].dropna())  # 欠損値を除去してソート
    cum_data = np.cumsum(data) / np.sum(data)  # 累積賃料比率を計算

    # グラフを描画
    plt.plot(data, cum_data, label=column)

plt.xlabel("賃料（円）")
plt.ylabel("累積賃料比率")
plt.title(f"賃料の累積比率グラフ ({datestamp})")  # タイトルに datestamp を追加
plt.legend()
plt.grid(True)

filename_cum1 = f"{datestamp}_cum1.png"  # ファイル名に datestamp を追加
image_path_cum1 = os.path.join(folder_path, filename_cum1)
plt.savefig(image_path_cum1)  # 正しいパスで保存
plt.close(fig)  # 画像を閉じる（メモリ解放のため）

print(f"画像が保存されました: \n{image_path_box1}\n{image_path_cum1}")

#まとめのパワポを作成する
# PowerPoint プレゼンテーションを作成
ppt = Presentation()

# スライドを追加
slide_layout1 = ppt.slide_layouts[0]  # タイトルスライド
slide1 = ppt.slides.add_slide(slide_layout1)

# タイトルを設定
title1 = slide1.shapes.title
title1.text = f"データサマリー ({datestamp})"
subtitle1 = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
subtitle1.text = "各駅の賃料をまとめました" 

# 📌【2ページ目】基礎統計量の表を追加
slide_layout = ppt.slide_layouts[5]  # タイトル＋コンテンツ
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "基礎統計量(小数桁数ご容赦)"

# `result_df.describe()` の統計データを取得
stats_df = result_df.describe()

# PowerPoint に表を追加
rows, cols = stats_df.shape
table = slide.shapes.add_table(rows+1, cols+1, Inches(1), Inches(1.5), Inches(8), Inches(4)).table

# ヘッダー行を挿入
table.cell(0, 0).text = "統計項目"
for col_idx, col_name in enumerate(stats_df.columns):
    table.cell(0, col_idx+1).text = col_name

# データ行を挿入
for row_idx, (index, row_data) in enumerate(stats_df.iterrows()):
    table.cell(row_idx+1, 0).text = index
    for col_idx, value in enumerate(row_data):
        table.cell(row_idx+1, col_idx+1).text = f"{value:.2f}"  # 小数点2桁に整形

# 📌【3ページ目】"賃料の箱ひげ図" + `filename_box1` の画像
slide_layout3 = ppt.slide_layouts[1]  # スライドマスター1（タイトルのみ）
slide3 = ppt.slides.add_slide(slide_layout)
slide3.shapes.title.text = "賃料の箱ひげ図"

# 画像を追加
image_path_box1 = os.path.join(folder_path, f"{datestamp}_box1.png")
if os.path.exists(image_path_box1):
    slide3.shapes.add_picture(image_path_box1, Inches(1), Inches(2), Inches(8), Inches(5))

# スライドにテキストを追加
text_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(1.2))  # 高さを広げる
text_frame = text_box.text_frame
text_frame.text = text_anova1  # ANOVAの結果を挿入

# テキストの折り返しを有効化
text_frame.word_wrap = True  

# フォントサイズを自動調整
for para in text_frame.paragraphs:
    para.font.size = Inches(0.2)  # 適切なフォントサイズに設定



# 📌【4ページ目】"賃料の分布" + `filename_cum1` の画像
slide_layout4 = ppt.slide_layouts[1]  # スライドマスター1（タイトルのみ）
slide4 = ppt.slides.add_slide(slide_layout)
slide4.shapes.title.text = "賃料の分布"

# 画像を追加
image_path_cum1 = os.path.join(folder_path, f"{datestamp}_cum1.png")
if os.path.exists(image_path_cum1):
    slide4.shapes.add_picture(image_path_cum1, Inches(1), Inches(2), Inches(8), Inches(5))


# PowerPoint ファイルの保存
ppt_filename = f"1c_{datestamp}_sum.pptx"
ppt_path = os.path.join(folder_path, ppt_filename)
ppt.save(ppt_path)


#ファルダ作成とファイル移動
# 作業ディレクトリの設定（元のフォルダーを指定）
source_folder = folder_path

# 現在の日時を取得して "AYYYYMMDDHHMM" のフォルダー名を作成
timestamp = datetime.now().strftime("A%Y%m%d%H%M")
dest_folder = os.path.join(source_folder, timestamp)

# フォルダーを作成
os.makedirs(dest_folder, exist_ok=True)

# ファイル移動
for filename in os.listdir(source_folder):
    file_path = os.path.join(source_folder, filename)
    
    if os.path.isfile(file_path):  # ファイルのみ対象
        if datestamp in filename:  # ファイル名に datestamp が含まれているかチェック
            shutil.move(file_path, os.path.join(dest_folder, filename))
            print(f"移動: {filename} → {dest_folder}")

print(f"ファイルの移動が完了しました。移動先: {dest_folder}")



