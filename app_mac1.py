# 統合版サーバー - HTMLファイルの配信とスクレイピングAPIを同時に提供（修正版）

import os
import time
import requests
from datetime import datetime
from bs4 import BeautifulSoup
import re
import pandas as pd
import numpy as np
import shutil
import scipy.stats as stats
from scipy.stats import linregress
import statsmodels.api as sm
from statsmodels.stats.outliers_influence import variance_inflation_factor
import matplotlib
matplotlib.use('Agg')  # バックエンドを明示的に設定
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from flask import Flask, request, jsonify, send_from_directory
import gc  # ガベージコレクション用

# Flaskアプリの設定
app = Flask(__name__)

# CORS設定（重複を避けるため簡潔に）
@app.after_request
def after_request(response):
    response.headers['Access-Control-Allow-Origin'] = '*'
    response.headers['Access-Control-Allow-Headers'] = 'Content-Type,Authorization'
    response.headers['Access-Control-Allow-Methods'] = 'GET,PUT,POST,DELETE,OPTIONS'
    response.headers['Access-Control-Max-Age'] = '86400'
    return response

# HTMLファイルを配信するルート
@app.route('/')
def index():
    return send_from_directory('.', 'keikyuuLine2.html')

# JavaScriptファイルを配信するルート
@app.route('/script.js')
def script():
    return send_from_directory('.', 'script.js')

@app.route('/<path:filename>')
def serve_file(filename):
    return send_from_directory('.', filename)

# テスト用のルート
@app.route('/test')
def test():
    return jsonify({"message": "Flask server is running!", "status": "OK"})

def cleanup_matplotlib():
    """matplotlib のリソースをクリーンアップ（必要最小限）"""
    plt.close('all')  # 全ての図を閉じる
    # rcParamsはリセットしない（日本語フォント設定を保持）

def cleanup_memory():
    """メモリクリーンアップ"""
    gc.collect()  # ガベージコレクション実行

# スクレイピングAPI（全機能統合版）
@app.route('/receive', methods=['POST', 'OPTIONS'])
def receive_and_scrape_data():
    print(f"Request method: {request.method}")
    
    if request.method == 'OPTIONS':
        print("Handling OPTIONS request")
        response = jsonify({})
        return response
    
    # 処理開始前にクリーンアップ
    cleanup_matplotlib()
    cleanup_memory()
    
    try:
        request_data = request.get_json()
        print(f"Received data: {request_data}")
    except Exception as e:
        print(f"JSON parse error: {e}")
        cleanup_matplotlib()
        cleanup_memory()
        return jsonify({"error": f"JSON parse error: {str(e)}"}), 400
    
    if not request_data:
        cleanup_matplotlib()
        cleanup_memory()
        return jsonify({"error": "Invalid JSON or no data"}), 400

    # JavaScriptから送られてきた値を処理
    email = request_data.get('email', '')
    num_pages = int(request_data.get('page', 3))
    stations = request_data.get('stations', ["三軒茶屋"])
    urls = request_data.get('urls', [
        "https://suumo.jp/jj/chintai/ichiran/FR301FC005/?ar=030&bs=040&ra=013&rn=0230&ek=023016720&cb=0.0&ct=9999999&mb=0&mt=9999999&et=9999999&cn=9999999&shkr1=03&shkr2=03&shkr3=03&shkr4=03&sngz=&po1=25&po2=99&pc=100&page="
    ])

    print(f"スクレイピング開始:")
    print(f"  Page数: {num_pages}")
    print(f"  駅リスト: {stations}")
    print(f"  Email: {email}")

    # フォルダーのパス
    folder_path = "/Users/su/Desktop/URLGen/1111accommodation"
    os.makedirs(folder_path, exist_ok=True)
    
    total_scraped = 0
    all_station_data = []  # 全駅のデータを保存するリスト
    
    try:
        # 各駅についてスクレイピング実行
        for station, base_url in zip(stations, urls):
            print(f"\n=== {station} のスクレイピング開始 ===")
            
            # 駅ごとの処理前にクリーンアップ
            cleanup_matplotlib()
            
            all_dataframes = []
            
            # 各ページをスクレイピング
            for i in range(1, num_pages + 1):
                url = base_url + str(i)
                print(f"📡 取得中: {url}")
                
                time.sleep(0.5)  # サイト負荷軽減
                
                try:
                    response = requests.get(url, timeout=10)
                    if response.status_code != 200:
                        print(f"ページ {i} の取得失敗: ステータスコード {response.status_code}")
                        continue
                except requests.exceptions.RequestException as e:
                    print(f"ページ {i} の取得中にエラー発生: {e}")
                    continue
                
                soup = BeautifulSoup(response.text, "html.parser")
            
                # 物件名・URLの取得
                titles = [title.text.strip() for title in soup.find_all("h2", class_="property_inner-title")]
                links = [a["href"] for a in soup.find_all("a", href=True) if "/chintai/bc" in a["href"]]
                full_links = ["https://suumo.jp" + link for link in links][:100]
            
                # 賃料の取得
                prices = [title.text.strip() for title in soup.find_all("div", class_="detailbox-property-point")]
                
                def convert_price(price):
                    try:
                        return int(float(price.replace('万円', '')) * 10000)
                    except ValueError:
                        return np.nan
                
                # 価格を変換してリストに格納
                rents = [convert_price(price) for price in prices]
            
                # 徒歩時間の取得
                walk_times = []
                detail_notes = soup.find_all("div", class_="font-weight:bold") + soup.find_all("div", style="font-weight:bold")
                
                for note in detail_notes:
                    text = note.text.strip()
                    try:
                        if "歩" in text and "分" in text and "バス" not in text and "車" not in text:
                            walk_time_str = text.split("歩")[1].split("分")[0].strip()
                            walk_time = int(walk_time_str)
                            walk_times.append(walk_time)
                        else:
                            walk_times.append(None)
                    except (ValueError, IndexError) as e:
                        print(f"⚠️ 変換できないデータ: {text}, エラー: {e}")
                        walk_times.append(None)
                
                # データフレームの長さを合わせる
                min_len = min(len(titles), len(full_links), len(rents), len(walk_times))
                
                if min_len == 0:
                    print(f"ページ {i}: データが見つかりませんでした")
                    continue
                
                # 物件情報 DataFrame
                df1 = pd.DataFrame({
                    "物件名": titles[:min_len],
                    "URL": full_links[:min_len],
                    "賃料(円)": rents[:min_len],
                    "徒歩時間(分)": walk_times[:min_len]
                })
                
                # 詳細情報取得
                properties = []
                for row in soup.find_all("tr")[:100]:  # 100件まで取得
                    try:
                        property_data = {
                            "管理費": row.find("td", class_="detailbox-property-col detailbox-property--col1").find_all("div")[1].text.strip(),
                            "敷金": row.find("td", class_="detailbox-property-col detailbox-property--col2").find_all("div")[0].text.strip(),
                            "礼金": row.find("td", class_="detailbox-property-col detailbox-property--col2").find_all("div")[1].text.strip(),
                            "間取り": row.find("td", class_="detailbox-property-col detailbox-property--col3").find_all("div")[0].text.strip(),
                            "専有面積(㎡)": row.find("td", class_="detailbox-property-col detailbox-property--col3").find_all("div")[1].text.strip(),
                            "向き": row.find("td", class_="detailbox-property-col detailbox-property--col3").find_all("div")[2].text.strip(),
                            "物件種別": row.find_all("td", class_="detailbox-property-col detailbox-property--col3")[1].find_all("div")[0].text.strip(),
                            "築年数(年)": row.find_all("td", class_="detailbox-property-col detailbox-property--col3")[1].find_all("div")[1].text.strip(),
                            "住所": row.find_all("td", class_="detailbox-property-col")[-1].text.strip()
                        }
                        properties.append(property_data)
                    except:
                        continue
                
                if properties:
                    df2 = pd.DataFrame(properties)
                    
                    # 専有面積をfloat型に変換
                    df2["専有面積(㎡)"] = df2["専有面積(㎡)"].str.replace("m2", "").astype(float)
                    
                    # 築年数をint型に変換
                    df2["築年数(年)"] = pd.to_numeric(df2["築年数(年)"].str.replace("築", "").str.replace("年", "").str.replace("新築", "0"), errors="coerce").astype("Int64")
                    df2["築年数(年)"] = df2["築年数(年)"].fillna(0).astype(int)
                    
                    # データフレーム結合
                    df_combined = pd.concat([df1, df2], axis=1)
                    all_dataframes.append(df_combined)
            
            # 全ページのデータを結合
            if not all_dataframes:
                print(f"{station}: データが1件も取得できませんでした。")
                continue

            final_df = pd.concat(all_dataframes, ignore_index=True)
            
            # データクリーニング
            for col in ['物件名', '向き']:
                if col in final_df.columns:
                    final_df[col] = final_df[col].astype(str)

            df_sorted = final_df.sort_values(by="物件名", ascending=True)
            
            # 不要な行を削除
            df_sorted = df_sorted[~df_sorted["物件名"].str.contains("築", na=False)]
            df_sorted = df_sorted[~df_sorted["物件名"].str.contains("号室", na=False)]
            df_sorted = df_sorted[~df_sorted["向き"].str.contains("-", na=False)]
            
            # 重複削除（元のコードに合わせて）
            columns = ['賃料(円)', '管理費', '間取り', '専有面積(㎡)', '向き']
            df_sorted = df_sorted.loc[~df_sorted[columns].eq(df_sorted[columns].shift(-1)).all(axis=1)]
            df_sorted = df_sorted.reset_index(drop=True)
            
            # 間取りでソート後に重複削除
            df_sorted = df_sorted.sort_values(by="間取り", ascending=True)
            df_sorted = df_sorted.sort_values(by="物件名", ascending=True)
            df_sorted = df_sorted.loc[~df_sorted[columns].eq(df_sorted[columns].shift(-1)).all(axis=1)]
            
            # 向きでソート後に重複削除
            df_sorted = df_sorted.sort_values(by="向き", ascending=True)
            df_sorted = df_sorted.sort_values(by="物件名", ascending=True)
            df_sorted = df_sorted.loc[~df_sorted[columns].eq(df_sorted[columns].shift(-1)).all(axis=1)]
            
            # 賃料でソート後に重複削除
            df_sorted = df_sorted.sort_values(by="賃料(円)", ascending=True)
            df_sorted = df_sorted.sort_values(by="物件名", ascending=True)
            df_sorted = df_sorted.loc[~df_sorted[columns].eq(df_sorted[columns].shift(-1)).all(axis=1)]
            
            df_sorted = df_sorted.reset_index(drop=True)
            
            # NAを削除
            df_sorted = df_sorted.replace('', pd.NA).dropna()

            n = len(df_sorted)
            total_scraped += n
            
            # 時刻生成
            datestamp = datetime.now().strftime("%y%m%d")
            timestamp = datetime.now().strftime("%y%m%d%H%M")
            
            # CSVファイル保存
            file_name = f"1fData_{station}_{datestamp}.csv"
            full_path = os.path.join(folder_path, file_name)
            
            try:
                df_sorted.to_csv(full_path, index=False, encoding="utf-8-sig")
                print(f"{station}: スクレイピング完了 - {n}件保存")
                
                # 統計処理とグラフ作成を実行
                print(f"{station}: 統計処理・グラフ作成開始")
                
                # matplotlib設定を各駅処理時に設定
                plt.rcParams['font.family'] = 'Hiragino Sans'
                
                # 基礎統計処理
                current_time = datetime.now()
                print(station)
                
                df_base1 = np.array([
                    ["全データ数", "取得した現在時刻", "調査駅", "出典"],
                    [n, "day"+timestamp, station, "https://suumo.jp/jj/chintai"]
                ])
                
                print(df_base1)
                
                # 賃料(円)の統計データ
                avg_total_rents = round(df_sorted["賃料(円)"].mean(),2)
                medi_total_rents = round(df_sorted["賃料(円)"].median(),2)
                stdevs_total_rents = round(df_sorted["賃料(円)"].std(ddof=1),4)
                std_error_total_rents = round(df_sorted["賃料(円)"].std(ddof=1)/np.sqrt(len(df_sorted)),4)
                min_total_rents = df_sorted["賃料(円)"].min()
                max_total_rents = df_sorted["賃料(円)"].max()
                firstQ_total_rents = round(df_sorted["賃料(円)"].quantile(0.25),1)
                thirdQ_total_rents = round(df_sorted["賃料(円)"].quantile(0.75),1)
                kurt_total_rents = round(df_sorted["賃料(円)"].kurt(),2)
                skew_total_rents = round(df_sorted["賃料(円)"].skew(),2)
                
                # 徒歩時間(分)の統計データ
                df_sorted["徒歩時間(分)"] = df_sorted["徒歩時間(分)"].astype(float)
                
                avg_total_walk_times = round(df_sorted["徒歩時間(分)"].mean(), 2)
                medi_total_walk_times = round(df_sorted["徒歩時間(分)"].median(), 2)
                stdevs_total_walk_times = round(df_sorted["徒歩時間(分)"].std(ddof=1), 4)
                std_error_total_walk_times = round(df_sorted["徒歩時間(分)"].std(ddof=1) / np.sqrt(len(df_sorted)), 4)
                min_total_walk_times = df_sorted["徒歩時間(分)"].min()
                max_total_walk_times = df_sorted["徒歩時間(分)"].max()
                firstQ_total_walk_times = round(df_sorted["徒歩時間(分)"].quantile(0.25), 2)
                thirdQ_total_walk_times = round(df_sorted["徒歩時間(分)"].quantile(0.75), 2)
                kurt_total_walk_times = round(df_sorted["徒歩時間(分)"].kurt(), 2)
                skew_total_walk_times = round(df_sorted["徒歩時間(分)"].skew(), 2)
                
                # 専有面積(㎡)の統計データ
                avg_total_space = round(df_sorted["専有面積(㎡)"].mean(),2)
                medi_total_space = round(df_sorted["専有面積(㎡)"].median(),2)
                stdevs_total_space = round(df_sorted["専有面積(㎡)"].std(ddof=1),4)
                std_error_total_space = round(df_sorted["専有面積(㎡)"].std(ddof=1)/np.sqrt(len(df_sorted)),4)
                min_total_space = df_sorted["専有面積(㎡)"].min()
                max_total_space = df_sorted["専有面積(㎡)"].max()
                firstQ_total_space = round(df_sorted["専有面積(㎡)"].quantile(0.25),2)
                thirdQ_total_space = round(df_sorted["専有面積(㎡)"].quantile(0.75),2)
                kurt_total_space = round(df_sorted["専有面積(㎡)"].kurt(),2)
                skew_total_space = round(df_sorted["専有面積(㎡)"].skew(),2)
                
                # 築年数(年)の統計データ
                avg_total_ages = round(df_sorted["築年数(年)"].mean(),2)
                medi_total_ages = round(df_sorted["築年数(年)"].median(),2)
                stdevs_total_ages = round(df_sorted["築年数(年)"].std(ddof=1),4)
                std_error_total_ages = round(df_sorted["築年数(年)"].std(ddof=1)/np.sqrt(len(df_sorted)),4)
                min_total_ages = df_sorted["築年数(年)"].min()
                max_total_ages = df_sorted["築年数(年)"].max()
                firstQ_total_ages = round(df_sorted["築年数(年)"].quantile(0.25),1)
                thirdQ_total_ages = round(df_sorted["築年数(年)"].quantile(0.75),1)
                kurt_total_ages = round(df_sorted["築年数(年)"].kurt(),2)
                skew_total_ages = round(df_sorted["築年数(年)"].skew(),2)
                
                # 統計データの配列作成
                df_stats1 = np.array([["項目","平均","中央値","不変標準偏差","標準誤差","最小値","最大値","第一四分位","第三四分位","尖度","歪度"],
                     ["賃料(円)",avg_total_rents,medi_total_rents,stdevs_total_rents,std_error_total_rents,min_total_rents,max_total_rents,firstQ_total_rents,thirdQ_total_rents,kurt_total_rents,skew_total_rents],
                    ["徒歩時間(分)",avg_total_walk_times,medi_total_walk_times,stdevs_total_walk_times,std_error_total_walk_times,min_total_walk_times,max_total_walk_times,firstQ_total_walk_times,thirdQ_total_walk_times,kurt_total_walk_times,skew_total_walk_times],
                    ["専有面積(㎡)",avg_total_space,medi_total_space,stdevs_total_space,std_error_total_space,min_total_space,max_total_space,firstQ_total_space,thirdQ_total_space,kurt_total_space,skew_total_space],
                    ["築年数(年)",avg_total_ages,medi_total_ages,stdevs_total_ages,std_error_total_ages,min_total_ages,max_total_ages,firstQ_total_ages,thirdQ_total_ages,kurt_total_ages,skew_total_ages]])
                
                df_stats1 = df_stats1.T
                print(df_stats1)
                
                # 分割統計データ
                df_stats11 = np.array([["項目","平均","中央値","不変標準偏差","標準誤差"],
                     ["賃料(円)",avg_total_rents,medi_total_rents,stdevs_total_rents,std_error_total_rents],
                    ["徒歩時間(分)",avg_total_walk_times,medi_total_walk_times,stdevs_total_walk_times,std_error_total_walk_times],
                    ["専有面積(㎡)",avg_total_space,medi_total_space,stdevs_total_space,std_error_total_space],
                    ["築年数(年)",avg_total_ages,medi_total_ages,stdevs_total_ages,std_error_total_ages]])
                
                df_stats11 = df_stats11.T
                
                df_stats12 = np.array([["項目","最小値","最大値","第一四分位","第三四分位","尖度","歪度"],
                     ["賃料(円)",min_total_rents,max_total_rents,firstQ_total_rents,thirdQ_total_rents,kurt_total_rents,skew_total_rents],
                    ["徒歩時間(分)",min_total_walk_times,max_total_walk_times,firstQ_total_walk_times,thirdQ_total_walk_times,kurt_total_walk_times,skew_total_walk_times],
                    ["専有面積(㎡)",min_total_space,max_total_space,firstQ_total_space,thirdQ_total_space,kurt_total_space,skew_total_space],
                    ["築年数(年)",min_total_ages,max_total_ages,firstQ_total_ages,thirdQ_total_ages,kurt_total_ages,skew_total_ages]])
                
                df_stats12 = df_stats12.T
                
                # 基本情報と基礎統計量をDataFrameに変換してCSV化
                df_base1 = pd.DataFrame(df_base1)
                df_stats1 = pd.DataFrame(df_stats1)
                
                file_name_base1 = f"{station}_{datestamp}_base1.csv"
                file_name_stats1 = f"{station}_{datestamp}_stats1.csv"
                
                full_path_base1 = os.path.join(folder_path, file_name_base1)
                full_path_stats1 = os.path.join(folder_path, file_name_stats1)
                
                df_base1.to_csv(full_path_base1, index=False, encoding="utf-8-sig")
                df_stats1.to_csv(full_path_stats1, index=False, encoding="utf-8-sig")
                
                # DataFrame分割化
                df_stats11 = pd.DataFrame(df_stats11)
                df_stats12 = pd.DataFrame(df_stats12)
                
                file_name_stats11 = f"{station}_{datestamp}_stats11.csv"
                file_name_stats12 = f"{station}_{datestamp}_stats12.csv"
                
                full_path_stats11 = os.path.join(folder_path, file_name_stats11)
                full_path_stats12 = os.path.join(folder_path, file_name_stats12)
                
                df_stats11.to_csv(full_path_stats11, index=False, encoding="utf-8-sig")
                df_stats12.to_csv(full_path_stats12, index=False, encoding="utf-8-sig")
                
                # グラフ作成（4×2レイアウト）
                fig, axes = plt.subplots(4, 2, figsize=(12, 12))
                
                # 賃料のヒストグラムと箱ひげ図
                axes[0, 0].hist(df_sorted["賃料(円)"], bins=30, edgecolor='black')
                axes[0, 0].set_title("賃料（円）のヒストグラム")
                axes[0, 0].set_xlabel("賃料（円）")
                axes[0, 0].set_ylabel("度数")
                
                axes[0, 1].boxplot(df_sorted["賃料(円)"], patch_artist=True, boxprops=dict(facecolor="skyblue"))
                axes[0, 1].set_title("賃料（円）の箱ひげ図")
                axes[0, 1].set_xlabel(station)
                axes[0, 1].set_ylabel("賃料（円）")
                
                # 徒歩時間のヒストグラムと箱ひげ図
                axes[1, 0].hist(df_sorted["徒歩時間(分)"], bins=30, edgecolor='black')
                axes[1, 0].set_title("徒歩時間（分）のヒストグラム")
                axes[1, 0].set_xlabel("徒歩時間（分）")
                axes[1, 0].set_ylabel("度数")
                
                axes[1, 1].boxplot(df_sorted["徒歩時間(分)"], patch_artist=True, boxprops=dict(facecolor="lightgreen"))
                axes[1, 1].set_title("徒歩時間（分）の箱ひげ図")
                axes[1, 1].set_xlabel(station)
                axes[1, 1].set_ylabel("徒歩時間（分）")
                
                # 専有面積のヒストグラムと箱ひげ図
                axes[2, 0].hist(df_sorted["専有面積(㎡)"], bins=30, edgecolor='black')
                axes[2, 0].set_title("専有面積(㎡)のヒストグラム")
                axes[2, 0].set_xlabel("専有面積(㎡)")
                axes[2, 0].set_ylabel("度数")
                
                axes[2, 1].boxplot(df_sorted["専有面積(㎡)"], patch_artist=True, boxprops=dict(facecolor="lightgreen"))
                axes[2, 1].set_title("専有面積(㎡)の箱ひげ図")
                axes[2, 1].set_xlabel(station)
                axes[2, 1].set_ylabel("専有面積(㎡)")
                
                # 築年数のヒストグラムと箱ひげ図
                axes[3, 0].hist(df_sorted["築年数(年)"], bins=30, edgecolor='black')
                axes[3, 0].set_title("築年数(年)のヒストグラム")
                axes[3, 0].set_xlabel("築年数(年)")
                axes[3, 0].set_ylabel("度数")
                
                axes[3, 1].boxplot(df_sorted["築年数(年)"], patch_artist=True, boxprops=dict(facecolor="lightgreen"))
                axes[3, 1].set_title("築年数(年)の箱ひげ図")
                axes[3, 1].set_xlabel(station)
                axes[3, 1].set_ylabel("築年数(年)")
                
                plt.tight_layout()
                file_name_g1 = f"{station}_{datestamp}_tg1.png"
                full_path_g1 = os.path.join(folder_path, file_name_g1)
                plt.savefig(full_path_g1)
                plt.close(fig)  # 明示的に図を閉じる
                
                # 個別グラフ作成
                # 賃料
                fig, axes = plt.subplots(1, 2, figsize=(12, 6))
                axes[0].hist(df_sorted["賃料(円)"], bins=30, edgecolor='black')
                axes[0].set_title("賃料(円)のヒストグラム")
                axes[0].set_xlabel("賃料(円)")
                axes[0].set_ylabel("度数")
                
                axes[1].boxplot(df_sorted["賃料(円)"], patch_artist=True, boxprops=dict(facecolor="skyblue"))
                axes[1].set_title("賃料(円)の箱ひげ図")
                axes[1].set_xlabel(station)
                axes[1].set_ylabel("賃料(円)")
                
                plt.tight_layout()
                file_name_gr1 = f"{station}_{datestamp}_gr1.png"
                full_path_gr1 = os.path.join(folder_path, file_name_gr1)
                plt.savefig(full_path_gr1)
                plt.close(fig)
                
                # 徒歩時間
                fig, axes = plt.subplots(1, 2, figsize=(12, 6))
                axes[0].hist(df_sorted["徒歩時間(分)"], bins=30, edgecolor='black')
                axes[0].set_title("徒歩時間(分)のヒストグラム")
                axes[0].set_xlabel("徒歩時間(分)")
                axes[0].set_ylabel("度数")
                
                axes[1].boxplot(df_sorted["徒歩時間(分)"], patch_artist=True, boxprops=dict(facecolor="skyblue"))
                axes[1].set_title("徒歩時間(分)の箱ひげ図")
                axes[1].set_xlabel(station)
                axes[1].set_ylabel("徒歩時間(分)")
                
                plt.tight_layout()
                file_name_gw1 = f"{station}_{datestamp}_gw1.png"
                full_path_gw1 = os.path.join(folder_path, file_name_gw1)
                plt.savefig(full_path_gw1)
                plt.close(fig)
                
                # 専有面積
                fig, axes = plt.subplots(1, 2, figsize=(12, 6))
                axes[0].hist(df_sorted["専有面積(㎡)"], bins=30, edgecolor='black')
                axes[0].set_title("専有面積のヒストグラム")
                axes[0].set_xlabel("専有面積(㎡)")
                axes[0].set_ylabel("度数")
                
                axes[1].boxplot(df_sorted["専有面積(㎡)"], patch_artist=True, boxprops=dict(facecolor="skyblue"))
                axes[1].set_title("専有面積の箱ひげ図")
                axes[1].set_xlabel(station)
                axes[1].set_ylabel("専有面積(㎡)")
                
                plt.tight_layout()
                file_name_gs1 = f"{station}_{datestamp}_gs1.png"
                full_path_gs1 = os.path.join(folder_path, file_name_gs1)
                plt.savefig(full_path_gs1)
                plt.close(fig)
                
                # 築年数
                fig, axes = plt.subplots(1, 2, figsize=(12, 6))
                axes[0].hist(df_sorted["築年数(年)"], bins=30, edgecolor='black')
                axes[0].set_title("築年数のヒストグラム")
                axes[0].set_xlabel("築年数(年)")
                axes[0].set_ylabel("度数")
                
                axes[1].boxplot(df_sorted["築年数(年)"], patch_artist=True, boxprops=dict(facecolor="skyblue"))
                axes[1].set_title("築年数の箱ひげ図")
                axes[1].set_xlabel(station)
                axes[1].set_ylabel("築年数(年)")
                
                plt.tight_layout()
                file_name_ga1 = f"{station}_{datestamp}_ga1.png"
                full_path_ga1 = os.path.join(folder_path, file_name_ga1)
                plt.savefig(full_path_ga1)
                plt.close(fig)
                
                # 散布図作成
                x1 = df_sorted["徒歩時間(分)"]
                x2 = df_sorted["専有面積(㎡)"]
                x3 = df_sorted["築年数(年)"]
                
                y1 = df_sorted["賃料(円)"]
                y2 = df_sorted["賃料(円)"]
                y3 = df_sorted["賃料(円)"]
                
                fig, axes = plt.subplots(3, 1, figsize=(8, 15))
                
                titles = ["賃料(円) vs 徒歩時間(分)", "賃料(円) vs 専有面積(㎡)", "賃料(円) vs 築年数(年)"]
                x_labels = ["徒歩時間(分)", "専有面積(㎡)", "築年数(年)"]
                y_labels = ["賃料(円)", "賃料(円)", "賃料(円)"]
                x_values = [x1, x2, x3]
                y_values = [y1, y2, y3]
                
                # グラフ描画
                for i in range(3):
                    # 線形回帰を計算
                    slope, intercept, r_value, p_value, std_err = linregress(x_values[i], y_values[i])
                    line_eq = f"y = {slope:.2f}x + {intercept:.2f}"
                    
                    # 散布図を描画
                    axes[i].scatter(x_values[i], y_values[i], alpha=0.6, color="blue", label="データ", s=20)
                    
                    # 近似直線を描画（x軸の範囲で）
                    x_line = np.linspace(x_values[i].min(), x_values[i].max(), 100)
                    y_line = slope * x_line + intercept
                    axes[i].plot(x_line, y_line, color="red", linewidth=2, label=f"近似直線: {line_eq}")
                    
                    # 決定係数とp値を表示（位置を調整）
                    x_pos = x_values[i].min() + (x_values[i].max() - x_values[i].min()) * 0.05
                    y_pos = y_values[i].max() - (y_values[i].max() - y_values[i].min()) * 0.15
                    axes[i].text(x_pos, y_pos, f"R² = {r_value**2:.3f}\np値 = {p_value:.4f}", 
                               fontsize=10, color="black", 
                               bbox=dict(boxstyle="round,pad=0.3", facecolor="white", alpha=0.8))
                    
                    # タイトル・ラベル設定
                    axes[i].set_title(titles[i], fontsize=12, pad=10)
                    axes[i].set_xlabel(x_labels[i], fontsize=10)
                    axes[i].set_ylabel(y_labels[i], fontsize=10)
                    axes[i].legend(loc="upper right")
                    axes[i].grid(True, alpha=0.3)
                
                plt.tight_layout()
                file_name_g2 = f"{station}_{datestamp}_tg2.png"
                full_path_g2 = os.path.join(folder_path, file_name_g2)
                plt.savefig(full_path_g2, dpi=200, bbox_inches='tight')
                plt.close(fig)
                
                # 散布図を個別に作成
                for i in range(3):
                    fig = plt.figure(figsize=(8, 6))
                    
                    # 線形回帰を計算
                    slope, intercept, r_value, p_value, std_err = linregress(x_values[i], y_values[i])
                    line_eq = f"y = {slope:.2f}x + {intercept:.2f}"
                    
                    # 散布図を描画
                    plt.scatter(x_values[i], y_values[i], alpha=0.6, color="blue", label="データ", s=30)
                    
                    # 近似直線を描画（正しい範囲で）
                    x_line = np.linspace(x_values[i].min(), x_values[i].max(), 100)
                    y_line = slope * x_line + intercept
                    plt.plot(x_line, y_line, color="red", linewidth=2, label=f"近似直線: {line_eq}")
                    
                    # 決定係数とp値を表示（位置を調整）
                    x_pos = x_values[i].min() + (x_values[i].max() - x_values[i].min()) * 0.05
                    y_pos = y_values[i].max() - (y_values[i].max() - y_values[i].min()) * 0.15
                    plt.text(x_pos, y_pos, f"R² = {r_value**2:.3f}\np値 = {p_value:.4f}", 
                            fontsize=11, color="black",
                            bbox=dict(boxstyle="round,pad=0.3", facecolor="white", alpha=0.8))
                    
                    # タイトル・ラベル設定
                    plt.title(titles[i], fontsize=14, pad=15)
                    plt.xlabel(x_labels[i], fontsize=12)
                    plt.ylabel(y_labels[i], fontsize=12)
                    plt.legend(loc="upper right")
                    plt.grid(True, alpha=0.3)
                    
                    # レイアウト調整
                    plt.tight_layout()
                    
                    # 画像の保存
                    file_name_tgscat = f"{station}_{datestamp}_tgscat{i+1}.png"
                    full_path_tgscat = os.path.join(folder_path, file_name_tgscat)
                    plt.savefig(full_path_tgscat, dpi=200, bbox_inches='tight')
                    plt.close(fig)
                
                # 間取り分類
                categories = ["ワンルーム", "1K", "1DK", "1LDK", "2K", "2DK", "2LDK", "3K", "3DK", "3LDK"]
                df_sorted["間取り分類"] = df_sorted["間取り"].apply(lambda x: x if x in categories else "その他")
                
                # グループ化して集計
                cat1 = df_sorted.groupby("間取り分類").agg(
                    件数=("間取り分類", "count"),
                    平均賃料=("賃料(円)", "mean"),
                    平均専有面積=("専有面積(㎡)", "mean")
                ).reset_index()
                
                cat1[["平均賃料", "平均専有面積"]] = cat1[["平均賃料", "平均専有面積"]].round(1)
                
                file_name_cat1 = f"{station}_{datestamp}_ct1.csv"
                full_path_cat1 = os.path.join(folder_path, file_name_cat1)
                cat1.to_csv(full_path_cat1, index=False, encoding="utf-8-sig")
                
                print(cat1)
                
                # 重回帰分析
                X = df_sorted[['徒歩時間(分)', '築年数(年)', '専有面積(㎡)']]
                y = df_sorted['賃料(円)']
                
                # 定数項を追加
                X = sm.add_constant(X)
                
                # 線形回帰モデルの作成
                model = sm.OLS(y, X).fit()
                
                print("   ")
                print(model.summary())
                print("-----切片を除いてP>|t|が0.05以下だと有意と考える　それ以上なら本来はその係数を除いて重回帰やり直し　関係性があるとは言い切れない-----")
                print("   ")
                
                adj_r_squared = model.rsquared_adj
                f_stat = model.fvalue
                f_p_value = model.f_pvalue
                intercept_coef = model.params["const"]
                coefficients = model.params.drop("const")
                p_values = model.pvalues.drop("const")
                
                # 日本語形式で表示
                print("=== 線形回帰モデルの結果 ===")
                print(f"補正決定係数: {adj_r_squared:.4f}")
                print(f"F値: {f_stat:.4f}")
                print(f"Fのp値: {f_p_value:.4f}")
                print(f"切片の係数: {intercept_coef:.4f}\n")
                
                print("各説明変数の傾きと p 値:")
                for var in coefficients.index:
                    print(f" - {var}: 傾き = {coefficients[var]:.4f}, p 値 = {p_values[var]:.4f}")
                
                # 重回帰基本情報のCSV保存
                df_mrl1 = np.array([
                    ["指標", "値"],
                    ["補正決定係数", adj_r_squared],
                    ["F値", f_stat],
                    ["Fのp値", f_p_value]
                ])
                
                df_mrl1 = pd.DataFrame(df_mrl1).T
                
                file_name_mrl1 = f"{station}_{datestamp}_mrl1.csv"
                full_path_mrl1 = os.path.join(folder_path, file_name_mrl1)
                df_mrl1.to_csv(full_path_mrl1, index=False, encoding="utf-8-sig")
                
                # 係数情報のCSV保存
                df_mrl2 = np.array([
                    ["item", "coef(切片、傾き)","p値"],
                    ["切片", intercept_coef,"-"],
                    ["徒歩時間(分)", coefficients["徒歩時間(分)"],p_values["徒歩時間(分)"]],
                    ["築年数(年)", coefficients["築年数(年)"],p_values["築年数(年)"]],
                    ["専有面積(㎡)", coefficients["専有面積(㎡)"],p_values["専有面積(㎡)"]]
                ])
                
                df_mrl2 = pd.DataFrame(df_mrl2)
                
                file_name_mrl2 = f"{station}_{datestamp}_mrl2.csv"
                full_path_mrl2 = os.path.join(folder_path, file_name_mrl2)
                df_mrl2.to_csv(full_path_mrl2, index=False, encoding="utf-8-sig")
                
                # 予測値と実測値の比較
                df_plot = df_sorted.copy()
                df_plot = df_plot.drop_duplicates()
                df_plot = df_plot.reset_index(drop=True)
                
                # 予測値を計算
                X_pred = sm.add_constant(df_plot[['徒歩時間(分)', '築年数(年)', '専有面積(㎡)']])
                df_plot['predicted_rent'] = model.predict(X_pred)
                
                # 残差の標準誤差を計算
                residuals = df_plot['賃料(円)'] - df_plot['predicted_rent']
                std_residuals = np.std(residuals)
                
                # 予測区間を残差の標準誤差で近似
                df_plot['upper_bound'] = df_plot['predicted_rent'] + (std_residuals * 1.96)
                df_plot['lower_bound'] = df_plot['predicted_rent'] - (std_residuals * 1.96)
                
                # プロット用に予測値でソート
                df_plot_sorted = df_plot.sort_values('predicted_rent').reset_index(drop=True)
                
                # 決定係数 (R²) の計算
                r_squared = model.rsquared
                p_values_model = model.pvalues
                
                # データ数を取得
                n_samples = len(df_plot)
                
                # 近似式の作成
                slope, intercept = np.polyfit(df_plot['賃料(円)'], df_plot['predicted_rent'], 1)
                line_eq = f"y = {slope:.2f}x + {intercept:.2f}"
                
                # 予測区間の幅を計算
                gap_pred = std_residuals * 1.96
                
                print(f"予測区間の幅（±1.96σ): {gap_pred:.1f}")
                
                # プロットの作成
                fig = plt.figure(figsize=(12, 8))
                
                # 散布図（実家賃 vs 予測家賃）
                plt.scatter(df_plot['賃料(円)'], df_plot['predicted_rent'], 
                           color="blue", alpha=0.6, label="実測値", s=30)
                
                # スムーズな線を描画するために十分な点を生成
                x_smooth = np.linspace(df_plot['賃料(円)'].min(), df_plot['賃料(円)'].max(), 100)
                y_smooth = slope * x_smooth + intercept
                
                # 回帰直線（スムーズ）
                plt.plot(x_smooth, y_smooth, "r-", lw=2, label="回帰直線")
                
                # 予測区間線をスムーズに描画
                upper_smooth = y_smooth + gap_pred
                lower_smooth = y_smooth - gap_pred
                
                plt.plot(x_smooth, upper_smooth, "k--", lw=1.5, alpha=0.8, label="予測区間上限")
                plt.plot(x_smooth, lower_smooth, "k--", lw=1.5, alpha=0.8, label="予測区間下限")
                
                # 予測区間の塗りつぶし
                plt.fill_between(x_smooth, lower_smooth, upper_smooth, 
                                 color="orange", alpha=0.2, label="予測区間")
                
                # 95%信頼区間も追加
                confidence_interval = std_residuals * 1.96 / np.sqrt(n_samples)
                upper_conf = y_smooth + confidence_interval
                lower_conf = y_smooth - confidence_interval
                
                plt.fill_between(x_smooth, lower_conf, upper_conf, 
                                 color="blue", alpha=0.3, label="95% 信頼区間")
                
                # グラフの詳細設定
                plt.xlabel("実際の賃料 (円)", fontsize=12)
                plt.ylabel("予測賃料 (円)", fontsize=12)
                plt.title("実際の賃料 vs 予測賃料（信頼区間・予測区間付き）", fontsize=14)
                plt.legend(loc='upper left')
                plt.grid(True, alpha=0.3)
                
                # 統計情報を右下に表示
                plt.text(0.98, 0.02, 
                         f"近似式: {line_eq}\nR² = {r_squared:.3f}\np値 = {p_values_model[1]:.3f}\nn = {n_samples}",
                         fontsize=11, verticalalignment="bottom", horizontalalignment="right",
                         transform=plt.gca().transAxes,
                         bbox=dict(facecolor="white", alpha=0.8, edgecolor="gray"))
                
                # 画像保存
                file_name_mlrap1 = f"{station}_{datestamp}_mlrap1.png"
                image_path_mlrap1 = os.path.join(folder_path, file_name_mlrap1)
                plt.savefig(image_path_mlrap1, dpi=300, bbox_inches='tight')
                plt.close(fig)
                
                # VIFの計算
                print("-----VIFは多重共線性（マルチコ；リニアリティー）の指数で1に近ければ大体OK-----")
                vif_data = pd.DataFrame()
                vif_data["feature"] = X.columns
                vif_data["VIF"] = [variance_inflation_factor(X.values, i) for i in range(X.shape[1])]
                
                # VIFをCSVに保存
                df_vif1 = np.array([
                    ["item", "VIF"],
                    ["徒歩時間(分)", variance_inflation_factor(X.values, 1)],
                    ["築年数(年)", variance_inflation_factor(X.values, 2)],
                    ["専有面積(㎡)", variance_inflation_factor(X.values, 3)]
                ])
                
                df_vif1 = pd.DataFrame(df_vif1)
                
                file_name_vif1 = f"{station}_{datestamp}_vif1.csv"
                full_path_vif1 = os.path.join(folder_path, file_name_vif1)
                df_vif1.to_csv(full_path_vif1, index=False, encoding="utf-8-sig")
                
                print(df_vif1)
                
                # 面積別の賃料予測値を計算
                pred25 = round(intercept_coef + coefficients["専有面積(㎡)"]*25 + coefficients["徒歩時間(分)"]*10 + coefficients["築年数(年)"]*10, 1)
                pred50 = round(intercept_coef + coefficients["専有面積(㎡)"]*50 + coefficients["徒歩時間(分)"]*10 + coefficients["築年数(年)"]*10, 1)
                pred75 = round(intercept_coef + coefficients["専有面積(㎡)"]*75 + coefficients["徒歩時間(分)"]*10 + coefficients["築年数(年)"]*10, 1)
                pred100 = round(intercept_coef + coefficients["専有面積(㎡)"]*100 + coefficients["徒歩時間(分)"]*10 + coefficients["築年数(年)"]*10, 1)
                
                # DataFrame を作成
                df_comp1 = pd.DataFrame([
                    ["25m²", pred25, round(pred25 - gap_pred, 1), round(pred25 + gap_pred, 1)],
                    ["50m²", pred50, round(pred50 - gap_pred, 1), round(pred50 + gap_pred, 1)],
                    ["75m²", pred75, round(pred75 - gap_pred, 1), round(pred75 + gap_pred, 1)],
                    ["100m²", pred100, round(pred100 - gap_pred, 1), round(pred100 + gap_pred, 1)]
                ], columns=["専有面積", "予測値", "予測下限", "予測上限"])
                
                print(df_comp1)
                
                # CSVファイルとして保存
                file_name_comp1 = f"{station}_{datestamp}_comp1.csv"
                full_path_comp1 = os.path.join(folder_path, file_name_comp1)
                df_comp1.to_csv(full_path_comp1, index=False, encoding="utf-8-sig")
                
                # PowerPoint作成
                file_name_ppt = f"1e_{station}_{timestamp}_ptt1.pptx"
                file_path_ppt = os.path.join(folder_path, file_name_ppt)
                
                # プレゼンテーションを作る
                pptt1 = Presentation()
                
                # タイトルスライドを追加
                slide_layout1 = pptt1.slide_layouts[0]
                slide1 = pptt1.slides.add_slide(slide_layout1)
                
                title = slide1.shapes.title
                subtitle = slide1.placeholders[1]
                
                title.text = f"{station}駅\n徒歩圏内の賃貸物件の\n調査結果"
                subtitle.text = f"調査時刻: {timestamp}\nデータ件数は{n}です\n ご注意:重複はなるべく排除していますが排除され切れていません"
                
                # 基本情報スライド
                slide_layout2 = pptt1.slide_layouts[6]
                slide2 = pptt1.slides.add_slide(slide_layout2)
                
                text_box2 = slide2.shapes.add_textbox(Cm(0.4), Cm(0.4), Cm(5), Cm(1))
                text_frame2 = text_box2.text_frame
                p2 = text_frame2.add_paragraph()
                p2.text = "基本情報"
                p2.font.size = Pt(16)
                p2.font.bold = True
                p2.font.color.rgb = RGBColor(0, 0, 0)
                
                # 表の作成
                rows, cols = df_base1.shape[0] + 1, df_base1.shape[1]
                table = slide2.shapes.add_table(rows, cols, Cm(1.5), Cm(2), Cm(22), Cm(4)).table
                
                # ヘッダー行の設定
                for col_idx, col_name in enumerate(df_base1.columns):
                    cell = table.cell(0, col_idx)
                    cell.text = str(col_name)
                
                # データ行の設定
                for row_idx, row in enumerate(df_base1.itertuples(), start=1):
                    for col_idx, value in enumerate(row[1:]):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(value)
                
                # カテゴリー情報スライド
                slide_layout3 = pptt1.slide_layouts[6]
                slide3 = pptt1.slides.add_slide(slide_layout3)
                
                text_box3 = slide3.shapes.add_textbox(Cm(0.4), Cm(0.4), Cm(5), Cm(1))
                text_frame3 = text_box3.text_frame
                p3 = text_frame3.add_paragraph()
                p3.text = "カテゴリー情報"
                p3.font.size = Pt(16)
                p3.font.bold = True
                p3.font.color.rgb = RGBColor(0, 0, 0)
                
                # cat1表の作成
                rows, cols = cat1.shape[0] + 1, cat1.shape[1]
                table = slide3.shapes.add_table(rows, cols, Cm(1.5), Cm(2), Cm(22), Cm(15)).table
                
                # ヘッダー行の設定
                for col_idx, col_name in enumerate(cat1.columns):
                    cell = table.cell(0, col_idx)
                    cell.text = str(col_name)
                
                # データ行の設定
                for row_idx, row in enumerate(cat1.itertuples(), start=1):
                    for col_idx, value in enumerate(row[1:]):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(value)
                
                # 基礎統計情報Aスライド
                slide_layout5 = pptt1.slide_layouts[6]
                slide5 = pptt1.slides.add_slide(slide_layout5)
                
                text_box5 = slide5.shapes.add_textbox(Cm(0.4), Cm(0.4), Cm(5), Cm(1))
                text_frame5 = text_box5.text_frame
                p5 = text_frame5.add_paragraph()
                p5.text = "基礎統計量情報A"
                p5.font.size = Pt(16)
                p5.font.bold = True
                p5.font.color.rgb = RGBColor(0, 0, 0)
                
                # df_stats11表の作成
                rows, cols = df_stats11.shape[0] + 1, df_stats11.shape[1]
                table = slide5.shapes.add_table(rows, cols, Cm(1.5), Cm(2), Cm(22), Cm(15)).table
                
                # ヘッダー行の設定
                for col_idx, col_name in enumerate(df_stats11.columns):
                    cell = table.cell(0, col_idx)
                    cell.text = str(col_name)
                
                # データ行の設定
                for row_idx, row in enumerate(df_stats11.itertuples(), start=1):
                    for col_idx, value in enumerate(row[1:]):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(value)
                
                # 基礎統計情報Bスライド
                slide_layout6 = pptt1.slide_layouts[6]
                slide6 = pptt1.slides.add_slide(slide_layout6)
                
                text_box6 = slide6.shapes.add_textbox(Cm(0.4), Cm(0.4), Cm(5), Cm(1))
                text_frame6 = text_box6.text_frame
                p6 = text_frame6.add_paragraph()
                p6.text = "基礎統計量情報B"
                p6.font.size = Pt(16)
                p6.font.bold = True
                p6.font.color.rgb = RGBColor(0, 0, 0)
                
                # df_stats12表の作成
                rows, cols = df_stats12.shape[0] + 1, df_stats12.shape[1]
                table = slide6.shapes.add_table(rows, cols, Cm(1.5), Cm(2), Cm(22), Cm(15)).table
                
                # ヘッダー行の設定
                for col_idx, col_name in enumerate(df_stats12.columns):
                    cell = table.cell(0, col_idx)
                    cell.text = str(col_name)
                
                # データ行の設定
                for row_idx, row in enumerate(df_stats12.itertuples(), start=1):
                    for col_idx, value in enumerate(row[1:]):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(value)
                
                # 分布と一次回帰のグラフスライド
                slide_layout7 = pptt1.slide_layouts[5]
                slide7 = pptt1.slides.add_slide(slide_layout7)
                
                title7 = slide7.shapes.title
                if title7:
                    title7.text = "全体の分布と一次回帰のグラフ"
                
                left = Inches(0.3)
                top = Inches(1.5)
                width = Inches(9.5)
                height = Inches(4.5)
                
                # 分布グラフを追加
                tg1_path = os.path.normpath(os.path.join(folder_path, f"{station}_{datestamp}_tg1.png"))
                if os.path.exists(tg1_path):
                    slide7.shapes.add_picture(tg1_path, left, top, width, height)
                    print(f"✅ 分布グラフを追加しました: {tg1_path}")
                
                # 散布図スライド（3枚に分割）
                # 散布図1: 賃料 vs 徒歩時間
                slide_layout8 = pptt1.slide_layouts[5]
                slide8 = pptt1.slides.add_slide(slide_layout8)
                
                title8 = slide8.shapes.title
                if title8:
                    title8.text = "賃料と徒歩時間の散布図"
                
                image_path_tgscat1 = os.path.normpath(os.path.join(folder_path, f"{station}_{datestamp}_tgscat1.png"))
                if os.path.exists(image_path_tgscat1):
                    slide8.shapes.add_picture(image_path_tgscat1, left, top, width, height)
                    print(f"✅ 散布図1を追加しました: {image_path_tgscat1}")
                
                # 散布図2: 賃料 vs 専有面積
                slide_layout8b = pptt1.slide_layouts[5]
                slide8b = pptt1.slides.add_slide(slide_layout8b)
                
                title8b = slide8b.shapes.title
                if title8b:
                    title8b.text = "賃料と専有面積の散布図"
                
                image_path_tgscat2 = os.path.normpath(os.path.join(folder_path, f"{station}_{datestamp}_tgscat2.png"))
                if os.path.exists(image_path_tgscat2):
                    slide8b.shapes.add_picture(image_path_tgscat2, left, top, width, height)
                    print(f"✅ 散布図2を追加しました: {image_path_tgscat2}")
                
                # 散布図3: 賃料 vs 築年数
                slide_layout8c = pptt1.slide_layouts[5]
                slide8c = pptt1.slides.add_slide(slide_layout8c)
                
                title8c = slide8c.shapes.title
                if title8c:
                    title8c.text = "賃料と築年数の散布図"
                
                image_path_tgscat3 = os.path.normpath(os.path.join(folder_path, f"{station}_{datestamp}_tgscat3.png"))
                if os.path.exists(image_path_tgscat3):
                    slide8c.shapes.add_picture(image_path_tgscat3, left, top, width, height)
                    print(f"✅ 散布図3を追加しました: {image_path_tgscat3}")
                
                # 賃料分布グラフスライド
                slide_layout9 = pptt1.slide_layouts[5]
                slide9 = pptt1.slides.add_slide(slide_layout9)
                
                title9 = slide9.shapes.title
                if title9:
                    title9.text = "賃料分布グラフ"
                
                image_path_gr1 = os.path.normpath(os.path.join(folder_path, f"{station}_{datestamp}_gr1.png"))
                if os.path.exists(image_path_gr1):
                    slide9.shapes.add_picture(image_path_gr1, left, top, width, height)
                    print(f"✅ 賃料分布グラフを追加しました: {image_path_gr1}")
                
                # 徒歩時間グラフスライド
                slide_layout10 = pptt1.slide_layouts[5]
                slide10 = pptt1.slides.add_slide(slide_layout10)
                
                title10 = slide10.shapes.title
                if title10:
                    title10.text = "徒歩時間グラフ"
                
                image_path_gw1 = os.path.normpath(os.path.join(folder_path, f"{station}_{datestamp}_gw1.png"))
                if os.path.exists(image_path_gw1):
                    slide10.shapes.add_picture(image_path_gw1, left, top, width, height)
                    print(f"✅ 徒歩時間グラフを追加しました: {image_path_gw1}")
                
                # 専有面積グラフスライド
                slide_layout11 = pptt1.slide_layouts[5]
                slide11 = pptt1.slides.add_slide(slide_layout11)
                
                title11 = slide11.shapes.title
                if title11:
                    title11.text = "専有面積グラフ"
                
                image_path_gs1 = os.path.normpath(os.path.join(folder_path, f"{station}_{datestamp}_gs1.png"))
                if os.path.exists(image_path_gs1):
                    slide11.shapes.add_picture(image_path_gs1, left, top, width, height)
                    print(f"✅ 専有面積グラフを追加しました: {image_path_gs1}")
                
                # 築年数グラフスライド
                slide_layout12 = pptt1.slide_layouts[5]
                slide12 = pptt1.slides.add_slide(slide_layout12)
                
                title12 = slide12.shapes.title
                if title12:
                    title12.text = "築年数グラフ"
                
                image_path_ga1 = os.path.normpath(os.path.join(folder_path, f"{station}_{datestamp}_ga1.png"))
                if os.path.exists(image_path_ga1):
                    slide12.shapes.add_picture(image_path_ga1, left, top, width, height)
                    print(f"✅ 築年数グラフを追加しました: {image_path_ga1}")
                
                # 重回帰分析結果スライド
                slide_layout15 = pptt1.slide_layouts[5]
                slide15 = pptt1.slides.add_slide(slide_layout15)
                
                if slide15.shapes.title:
                    slide15.shapes.title.text = "重回帰分析結果"
                
                text_box15 = slide15.shapes.add_textbox(Cm(0.4), Cm(2), Cm(5), Cm(1))
                text_frame15 = text_box15.text_frame
                p15 = text_frame15.add_paragraph()
                p15.text = "重回帰基礎結果とcoefficients"
                p15.font.size = Pt(16)
                p15.font.bold = True
                p15.font.color.rgb = RGBColor(0, 0, 0)
                
                # table15_1をスライドの上から4cmの位置に配置
                table15_1 = slide15.shapes.add_table(df_mrl1.shape[0], df_mrl1.shape[1], Cm(1.5), Cm(4.0), Cm(22), Cm(4)).table
                
                # table15_2をスライドの上から10cmの位置に配置
                table15_2 = slide15.shapes.add_table(df_mrl2.shape[0], df_mrl2.shape[1], Cm(1.5), Cm(10.0), Cm(22), Cm(4)).table
                
                # table15_1にdf_mrl1のデータを追加
                for row_idx, (index, row) in enumerate(df_mrl1.iterrows()):
                    for col_idx, value in enumerate(row):
                        cell = table15_1.cell(row_idx, col_idx)
                        cell.text = str(value)
                
                # table15_2にdf_mrl2のデータを追加
                for row_idx, (index, row) in enumerate(df_mrl2.iterrows()):
                    for col_idx, value in enumerate(row):
                        cell = table15_2.cell(row_idx, col_idx)
                        cell.text = str(value)
                
                # VIFスライド
                slide_layout16 = pptt1.slide_layouts[6]
                slide16 = pptt1.slides.add_slide(slide_layout16)
                
                text_box16 = slide16.shapes.add_textbox(Cm(0.4), Cm(0.5), Cm(5), Cm(1))
                text_frame16 = text_box16.text_frame
                p16 = text_frame16.add_paragraph()
                p16.text = "重回帰の多重共線性（VIF)"
                p16.font.size = Pt(16)
                p16.font.bold = True
                p16.font.color.rgb = RGBColor(0, 0, 0)
                
                # df_vif1表の作成
                rows, cols = df_vif1.shape[0] + 1, df_vif1.shape[1]
                table = slide16.shapes.add_table(rows, cols, Cm(1.5), Cm(2.5), Cm(22), Cm(4)).table
                
                # ヘッダー行の設定
                for col_idx, col_name in enumerate(df_vif1.columns):
                    cell = table.cell(0, col_idx)
                    cell.text = str(col_name)
                
                # データ行の設定
                for row_idx, row in enumerate(df_vif1.itertuples(), start=1):
                    for col_idx, value in enumerate(row[1:]):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(value)
                
                # 予測家賃と実家賃の関係スライド
                slide_layout17 = pptt1.slide_layouts[5]
                slide17 = pptt1.slides.add_slide(slide_layout17)
                
                title17 = slide17.shapes.title
                if title17:
                    title17.text = "予測家賃と実家賃の関係"
                
                if os.path.exists(image_path_mlrap1):
                    slide17.shapes.add_picture(image_path_mlrap1, left, top, width, Inches(5.0))
                    print(f"✅ 予測vs実測グラフを追加しました: {image_path_mlrap1}")
                
                # 面積毎の家賃予測スライド
                slide_layout18 = pptt1.slide_layouts[6]
                slide18 = pptt1.slides.add_slide(slide_layout18)
                
                text_box18 = slide18.shapes.add_textbox(Cm(0.4), Cm(0.5), Cm(5), Cm(1))
                text_frame18 = text_box18.text_frame
                p18 = text_frame18.add_paragraph()
                p18.text = "面積毎の家賃予測"
                p18.font.size = Pt(16)
                p18.font.bold = True
                p18.font.color.rgb = RGBColor(0, 0, 0)
                
                # df_comp1表の作成
                rows, cols = df_comp1.shape[0] + 1, df_comp1.shape[1]
                table = slide18.shapes.add_table(rows, cols, Cm(1.5), Cm(2), Cm(22), Cm(4)).table
                
                # ヘッダー行の設定
                for col_idx, col_name in enumerate(df_comp1.columns):
                    cell = table.cell(0, col_idx)
                    cell.text = str(col_name)
                
                # データ行の設定
                for row_idx, row in enumerate(df_comp1.itertuples(), start=1):
                    for col_idx, value in enumerate(row[1:]):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(value)
                
                # フッター追加
                for index, slide in enumerate(pptt1.slides):
                    current_page = index + 1
                    total_pages = len(pptt1.slides)
                    
                    left_text = f"{station}, n={n}"
                    center_text = f"{current_page}/{total_pages}"
                    right_text = f"{timestamp}"
                    
                    # テキストボックスの追加（スライド下部）
                    left_box = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(2), Inches(0.3))
                    left_box.text_frame.text = left_text
                    
                    center_box = slide.shapes.add_textbox(Inches(4.2), Inches(7.15), Inches(2), Inches(0.3))
                    center_box.text_frame.text = center_text
                    
                    right_box = slide.shapes.add_textbox(Inches(8.0), Inches(7.15), Inches(2), Inches(0.3))
                    right_box.text_frame.text = right_text
                
                # PowerPointを保存
                pptt1.save(file_path_ppt)
                print(f"✅ PowerPointファイルを保存しました: {file_path_ppt}")
                
                # 駅データを全体リストに追加
                all_station_data.append({
                    'station': station,
                    'data': df_sorted,
                    'count': n
                })
                
                print(f"{station}: 統計処理・グラフ作成・PowerPoint作成完了")
                
            except Exception as e:
                print(f"{station}: 処理中にエラーが発生しました: {e}")
            finally:
                # 駅ごとの処理後にクリーンアップ
                cleanup_matplotlib()
                cleanup_memory()
        
        # 全駅のデータが揃った後、総合まとめを作成
        if all_station_data:
            print("\n=== 総合まとめ作成開始 ===")
            
            try:
                # 総合まとめの処理
                csv_files = [f for f in os.listdir(folder_path) if f.startswith("1fData") and f.endswith(".csv")]
                
                # データを格納するリスト
                data_list = []
                
                # CSVファイルを処理
                for file in csv_files:
                    file_path_csv = os.path.join(folder_path, file)
                    
                    # ファイル名から駅名を抽出
                    name_parts = file.split("_")
                    if len(name_parts) >= 3:
                        column_name = name_parts[1]
                        
                        # CSVを読み込む
                        df = pd.read_csv(file_path_csv)
                        
                        # 「賃料（円）」の列を抽出
                        rent_column = [col for col in df.columns if "賃料" in col or "円" in col]
                        if rent_column:
                            df_filtered = df[[rent_column[0]]]
                            df_filtered.columns = [column_name]
                            data_list.append(df_filtered)
                
                # 複数のデータを結合
                if data_list:
                    result_df = pd.concat(data_list, axis=1)
                    print(result_df)
                    
                    # 基礎統計量を求める
                    stats_df = result_df.describe()
                    print("基礎統計量:\n", stats_df)
                    
                    # matplotlib設定
                    plt.rcParams['font.family'] = 'Hiragino Sans'
                    
                    # 箱ひげ図の作成・保存
                    fig = plt.figure(figsize=(10, 6))
                    result_df.boxplot()
                    plt.title(f"箱ひげ図 ({datestamp})")
                    plt.ylabel("賃料（円）")
                    plt.xticks(rotation=45)
                    plt.grid(True)
                    
                    filename_box1 = f"{datestamp}_box1.png"
                    image_path_box1 = os.path.join(folder_path, filename_box1)
                    plt.savefig(image_path_box1)
                    plt.close(fig)
                    
                    # 列名を取得
                    column_names = result_df.columns.tolist()
                    
                    # 欠損値を削除してデータを整理
                    result_df = result_df.dropna()
                    
                    # ANOVAの実施
                    groups = [result_df[col] for col in column_names]
                    groups = [g for g in groups if len(g) > 0]
                    
                    if len(groups) > 1:
                        F_value, p_value = stats.f_oneway(*groups)
                        text_anova1 = f"一元配置分散分析（ANOVA）の結果:\nF値: {F_value:.2f}\np値: {p_value:.3f}"
                        print(text_anova1)
                    else:
                        text_anova1 = "ANOVAを適用できる十分なデータがありません。"
                        print(text_anova1)
                    
                    # 賃料の累積比率グラフの作成・保存
                    fig = plt.figure(figsize=(10, 6))
                    
                    for column in result_df.columns:
                        data = np.sort(result_df[column].dropna())
                        cum_data = np.cumsum(data) / np.sum(data)
                        plt.plot(data, cum_data, label=column)
                    
                    plt.xlabel("賃料（円）")
                    plt.ylabel("累積賃料比率")
                    plt.title(f"賃料の累積比率グラフ ({datestamp})")
                    plt.legend()
                    plt.grid(True)
                    
                    filename_cum1 = f"{datestamp}_cum1.png"
                    image_path_cum1 = os.path.join(folder_path, filename_cum1)
                    plt.savefig(image_path_cum1)
                    plt.close(fig)
                    
                    print(f"画像が保存されました: \n{image_path_box1}\n{image_path_cum1}")
                    
                    # まとめのパワポを作成する
                    ppt = Presentation()
                    
                    # スライドを追加
                    slide_layout1 = ppt.slide_layouts[0]
                    slide1 = ppt.slides.add_slide(slide_layout1)
                    
                    # タイトルを設定
                    title1 = slide1.shapes.title
                    title1.text = f"データサマリー ({datestamp})"
                    subtitle1 = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
                    subtitle1.text = "各駅の賃料をまとめました"
                    
                    # 2ページ目：基礎統計量の表を追加
                    slide_layout = ppt.slide_layouts[5]
                    slide = ppt.slides.add_slide(slide_layout)
                    title = slide.shapes.title
                    title.text = "基礎統計量(小数桁数ご容赦)"
                    
                    # result_df.describe()の統計データを取得
                    stats_df = result_df.describe()
                    
                    # PowerPointに表を追加
                    rows, cols = stats_df.shape
                    table = slide.shapes.add_table(rows+1, cols+1, Inches(1), Inches(1.5), Inches(8), Inches(4)).table
                    
                    # ヘッダー行を挿入
                    table.cell(0, 0).text = "統計項目"
                    for col_idx, col_name in enumerate(stats_df.columns):
                        table.cell(0, col_idx+1).text = col_name
                    
                    # データ行を挿入
                    for row_idx, (index, row_data) in enumerate(stats_df.iterrows()):
                        table.cell(row_idx+1, 0).text = index
                        for col_idx, value in enumerate(row_data):
                            table.cell(row_idx+1, col_idx+1).text = f"{value:.2f}"
                    
                    # 3ページ目："賃料の箱ひげ図"
                    slide_layout3 = ppt.slide_layouts[1]
                    slide3 = ppt.slides.add_slide(slide_layout3)
                    slide3.shapes.title.text = "賃料の箱ひげ図"
                    
                    # 画像を追加
                    if os.path.exists(image_path_box1):
                        slide3.shapes.add_picture(image_path_box1, Inches(1), Inches(2), Inches(8), Inches(5))
                    
                    # スライドにテキストを追加
                    text_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(1.2))
                    text_frame = text_box.text_frame
                    text_frame.text = text_anova1
                    text_frame.word_wrap = True
                    
                    for para in text_frame.paragraphs:
                        para.font.size = Inches(0.2)
                    
                    # 4ページ目："賃料の分布"
                    slide_layout4 = ppt.slide_layouts[1]
                    slide4 = ppt.slides.add_slide(slide_layout4)
                    slide4.shapes.title.text = "賃料の分布"
                    
                    # 画像を追加
                    if os.path.exists(image_path_cum1):
                        slide4.shapes.add_picture(image_path_cum1, Inches(1), Inches(2), Inches(8), Inches(5))
                    
                    # PowerPointファイルの保存
                    ppt_filename = f"1c_{datestamp}_sum.pptx"
                    ppt_path = os.path.join(folder_path, ppt_filename)
                    ppt.save(ppt_path)
                    print(f"✅ 総合まとめPowerPointを保存しました: {ppt_path}")
                else:
                    print("CSVファイルが見つかりませんでした")
                
            except Exception as e:
                print(f"総合まとめ作成中にエラーが発生しました: {e}")
            finally:
                cleanup_matplotlib()
                cleanup_memory()
            
            print("=== 総合まとめ作成完了 ===")
        
        # ファイル移動
        print("\n=== ファイル移動開始 ===")
        
        try:
            # 現在の日時を取得して "AYYYYMMDDHHMM" のフォルダー名を作成
            timestamp_folder = datetime.now().strftime("A%Y%m%d%H%M")
            dest_folder = os.path.join(folder_path, timestamp_folder)
            
            # フォルダーを作成
            os.makedirs(dest_folder, exist_ok=True)
            
            # ファイル移動
            for filename in os.listdir(folder_path):
                file_path = os.path.join(folder_path, filename)
                
                if os.path.isfile(file_path):
                    if datestamp in filename:
                        shutil.move(file_path, os.path.join(dest_folder, filename))
                        print(f"移動: {filename} → {dest_folder}")
            
            print(f"ファイルの移動が完了しました。移動先: {dest_folder}")
            
        except Exception as e:
            print(f"ファイル移動中にエラーが発生しました: {e}")
            dest_folder = folder_path
        
    except Exception as e:
        print(f"処理中に予期しないエラーが発生しました: {e}")
        cleanup_matplotlib()
        cleanup_memory()
        return jsonify({"error": f"処理中にエラーが発生しました: {str(e)}"}), 500
    
    finally:
        # 最終クリーンアップ
        cleanup_matplotlib()
        cleanup_memory()
    
    return jsonify({
        "message": "全処理完了！スクレイピング、統計処理、グラフ作成、PowerPoint作成、ファイル移動が完了しました。",
        "page": num_pages,
        "email": email,
        "stations": stations,
        "scraped_stations": len(stations),
        "total_scraped_items": total_scraped,
        "output_folder": dest_folder if 'dest_folder' in locals() else folder_path
    })

if __name__ == '__main__':
    print("=" * 50)
    print("修正版完全統合サーバー起動中...")
    print("=" * 50)
    print("1. HTMLファイル 'keikyuuLine2.html' をこのPythonファイルと同じフォルダに置いてください")
    print("2. ブラウザで http://localhost:5000 にアクセスしてください")
    print("3. テスト用URL: http://localhost:5000/test")
    print("4. 機能: スクレイピング → 統計処理 → グラフ作成 → PowerPoint作成 → ファイル移動")
    print("5. 修正点: matplotlib/メモリクリーンアップ対応で再起動不要")
    print("=" * 50)
    
    app.run(debug=True, host='0.0.0.0', port=5000)